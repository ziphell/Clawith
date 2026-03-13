"""Agent tools — unified file-based tools that give digital employees
access to their own structured workspace.

Design principle:  ONE set of file tools covers EVERYTHING.
The agent's workspace uses well-known paths:
  - tasks.json          → task list (auto-synced from DB)
  - soul.md             → personality definition
  - memory.md           → long-term memory / notes
  - skills/             → skill definitions (markdown files)
  - workspace/          → general working files, reports, etc.

The agent reads/writes these files directly. No per-concept tools needed.
"""

import json
import os
import uuid
from contextvars import ContextVar
from datetime import datetime, timezone
from pathlib import Path

from sqlalchemy import select

from app.database import async_session
from app.models.task import Task
from app.config import get_settings

_settings = get_settings()
WORKSPACE_ROOT = Path(_settings.AGENT_DATA_DIR)

# ContextVar set by each channel handler so send_channel_file knows where to send
# Value: async callable(file_path: Path) -> None  |  None for web chat (returns URL)
channel_file_sender: ContextVar = ContextVar('channel_file_sender', default=None)
# For web chat: agent_id needed to build download URL
channel_web_agent_id: ContextVar = ContextVar('channel_web_agent_id', default=None)
# Set by Feishu channel handler — open_id of the message sender so calendar tool
# can auto-invite them as attendee when no explicit attendee list is given
channel_feishu_sender_open_id: ContextVar = ContextVar('channel_feishu_sender_open_id', default=None)

# ─── Tool Definitions (OpenAI function-calling format) ──────────

AGENT_TOOLS = [
    {
        "type": "function",
        "function": {
            "name": "list_files",
            "description": "List files and folders in a directory within my workspace. Can also list enterprise_info/ for shared company information.",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {
                        "type": "string",
                        "description": "Directory path to list, defaults to root (empty string). e.g.: '', 'skills', 'workspace', 'enterprise_info', 'enterprise_info/knowledge_base'",
                    }
                },
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "read_file",
            "description": "Read file contents from the workspace. Can read tasks.json for tasks, soul.md for personality, memory/memory.md for memory, skills/ for skill files, and enterprise_info/ for shared company info.",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {
                        "type": "string",
                        "description": "File path, e.g.: tasks.json, soul.md, memory/memory.md, skills/xxx.md, enterprise_info/company_profile.md",
                    }
                },
                "required": ["path"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "write_file",
            "description": "Write or update a file in the workspace. Can update memory/memory.md, focus.md, task_history.md, create documents in workspace/, create skills in skills/.",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {
                        "type": "string",
                        "description": "File path, e.g.: memory/memory.md, workspace/report.md, skills/data_analysis.md",
                    },
                    "content": {
                        "type": "string",
                        "description": "File content to write",
                    },
                },
                "required": ["path", "content"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "delete_file",
            "description": "Delete a file from the workspace. Cannot delete soul.md or tasks.json.",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {
                        "type": "string",
                        "description": "File path to delete",
                    }
                },
                "required": ["path"],
            },
        },
    },
    # --- Trigger management tools (Pulse engine) ---
    {
        "type": "function",
        "function": {
            "name": "set_trigger",
            "description": "Set a new trigger to wake yourself up at a specific time or condition. Use this to schedule future actions, monitor changes, or wait for messages. The trigger will fire and invoke you with the reason text as context. Trigger types: 'cron' (recurring schedule), 'once' (fire once at a time), 'interval' (every N minutes), 'poll' (HTTP monitoring), 'on_message' (when another agent or a human user replies — use from_agent_name for agents, or from_user_name for human users on Feishu/Slack/Discord), 'webhook' (receive external HTTP POST — system generates a unique URL, give it to the user so they can configure it in external services like GitHub, Grafana, etc.).",
            "parameters": {
                "type": "object",
                "properties": {
                    "name": {
                        "type": "string",
                        "description": "Unique name for this trigger, e.g. 'daily_briefing' or 'wait_morty_reply'",
                    },
                    "type": {
                        "type": "string",
                        "enum": ["cron", "once", "interval", "poll", "on_message", "webhook"],
                        "description": "Trigger type",
                    },
                    "config": {
                        "type": "object",
                        "description": "Type-specific config. cron: {\"expr\": \"0 9 * * *\"}. once: {\"at\": \"2026-03-10T09:00:00+08:00\"}. interval: {\"minutes\": 30}. poll: {\"url\": \"...\", \"json_path\": \"$.status\", \"fire_on\": \"change\", \"interval_min\": 5}. on_message: {\"from_agent_name\": \"Morty\"} or {\"from_user_name\": \"张三\"} (for human users on Feishu/Slack/Discord). webhook: {\"secret\": \"optional_hmac_secret\"} (system auto-generates the URL)",
                    },
                    "reason": {
                        "type": "string",
                        "description": "What you should do when this trigger fires. This will be shown to you as context when you wake up.",
                    },
                    "focus_ref": {
                        "type": "string",
                        "description": "Optional: identifier of the focus item in focus.md that this trigger relates to (use the checklist identifier, e.g. 'daily_news_check')",
                    },
                },
                "required": ["name", "type", "config", "reason"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "update_trigger",
            "description": "Update an existing trigger's configuration or reason. Use this to adjust timing, change parameters, etc. For example, change interval from 5 minutes to 30 minutes.",
            "parameters": {
                "type": "object",
                "properties": {
                    "name": {
                        "type": "string",
                        "description": "Name of the trigger to update",
                    },
                    "config": {
                        "type": "object",
                        "description": "New config (replaces existing config)",
                    },
                    "reason": {
                        "type": "string",
                        "description": "New reason text",
                    },
                },
                "required": ["name"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "cancel_trigger",
            "description": "Cancel (disable) a trigger by name. Use this when a task is completed and the trigger is no longer needed.",
            "parameters": {
                "type": "object",
                "properties": {
                    "name": {
                        "type": "string",
                        "description": "Name of the trigger to cancel",
                    },
                },
                "required": ["name"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "list_triggers",
            "description": "List all your active triggers. Shows name, type, config, reason, fire count, and status.",
            "parameters": {
                "type": "object",
                "properties": {},
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "send_channel_file",
            "description": "Send a file to the user via the current communication channel (Feishu, Slack, Discord, or web). Call this when you have created a file and the user would benefit from receiving it directly. Provide the workspace-relative file path (e.g. workspace/report.md).",
            "parameters": {
                "type": "object",
                "properties": {
                    "file_path": {
                        "type": "string",
                        "description": "Workspace-relative path to the file, e.g. workspace/report.md",
                    },
                    "message": {
                        "type": "string",
                        "description": "Optional message to accompany the file",
                    },
                },
                "required": ["file_path"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "send_feishu_message",
            "description": (
                "Send a Feishu IM message to a colleague. "
                "You can provide either the colleague's name (will auto-search their open_id) "
                "or their open_id directly. "
                "To contact digital employees use send_message_to_agent instead."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "member_name": {
                        "type": "string",
                        "description": "Recipient's name, e.g. '覃睿'. Will be looked up automatically.",
                    },
                    "open_id": {
                        "type": "string",
                        "description": "Recipient's Feishu open_id (e.g. from feishu_user_search). Use this if you already have it.",
                    },
                    "message": {
                        "type": "string",
                        "description": "Message content to send",
                    },
                },
                "required": ["message"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "send_web_message",
            "description": "Send a message to a user on the Clawith web platform. The message will appear in their web chat history and be pushed in real-time if they are online. Use this to proactively notify web users.",
            "parameters": {
                "type": "object",
                "properties": {
                    "username": {
                        "type": "string",
                        "description": "Username or display name of the recipient (must be a registered platform user)",
                    },
                    "message": {
                        "type": "string",
                        "description": "Message content to send",
                    },
                },
                "required": ["username", "message"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "send_message_to_agent",
            "description": "Send a message to a digital employee colleague and receive a reply. The recipient is another AI agent, not a human. This triggers the recipient's LLM reasoning and returns their response. Suitable for asking questions, delegating tasks, or collaboration. Your relationships.md lists available digital employees under 'Digital Employee Colleagues'.",
            "parameters": {
                "type": "object",
                "properties": {
                    "agent_name": {
                        "type": "string",
                        "description": "Target digital employee's name",
                    },
                    "message": {
                        "type": "string",
                        "description": "Message content to send",
                    },
                    "msg_type": {
                        "type": "string",
                        "enum": ["notify", "consult", "task_delegate"],
                        "description": "Message type: notify (notification), consult (ask a question), task_delegate (delegate a task). Defaults to notify.",
                    },
                },
                "required": ["agent_name", "message"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "jina_search",
            "description": "Search the internet using Jina AI Search (s.jina.ai). Returns high-quality search results with full page content, not just snippets. Ideal for research, news, technical docs, and any real-time information lookup.",
            "parameters": {
                "type": "object",
                "properties": {
                    "query": {
                        "type": "string",
                        "description": "Search query, e.g. 'Python asyncio best practices' or '苏州通道人工智能科技有限公司'",
                    },
                    "max_results": {
                        "type": "integer",
                        "description": "Number of results to return, default 5, max 10",
                    },
                },
                "required": ["query"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "jina_read",
            "description": "Read and extract the full content from a web page URL using Jina AI Reader (r.jina.ai). Returns clean, well-structured markdown including article text, tables, and key information. Better than jina_search when you already have a specific URL to read.",
            "parameters": {
                "type": "object",
                "properties": {
                    "url": {
                        "type": "string",
                        "description": "The full URL of the web page to read, e.g. 'https://example.com/article'",
                    },
                    "max_chars": {
                        "type": "integer",
                        "description": "Max characters to return (default 8000, max 20000)",
                    },
                },
                "required": ["url"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "read_document",
            "description": "Read office document contents (PDF, Word, Excel, PPT, etc.) and extract text. Suitable for reading knowledge base documents.",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {
                        "type": "string",
                        "description": "Document file path, e.g.: workspace/knowledge_base/report.pdf, enterprise_info/knowledge_base/policy.docx",
                    }
                },
                "required": ["path"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "execute_code",
            "description": "Execute code (Python, Bash, or Node.js) in a sandboxed environment within the agent's workspace directory. Useful for data processing, calculations, file transformations, and automation scripts. Code runs with the workspace as the working directory. Security restrictions apply: no network access commands, no system-level operations, 30-second timeout.",
            "parameters": {
                "type": "object",
                "properties": {
                    "language": {
                        "type": "string",
                        "enum": ["python", "bash", "node"],
                        "description": "Programming language to execute",
                    },
                    "code": {
                        "type": "string",
                        "description": "Code to execute. For Python, you can import standard libraries (json, csv, math, re, collections, etc.). Working directory is your workspace/.",
                    },
                    "timeout": {
                        "type": "integer",
                        "description": "Max execution time in seconds (default 30, max 60)",
                    },
                },
                "required": ["language", "code"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "upload_image",
            "description": "Upload an image file from your workspace (or from a public URL) to a cloud CDN and get a permanent public URL. Use this when you need to share images externally, embed them in messages/reports, or make workspace images accessible via URL. Supports common formats: PNG, JPG, GIF, WebP, SVG.",
            "parameters": {
                "type": "object",
                "properties": {
                    "file_path": {
                        "type": "string",
                        "description": "Workspace-relative path to the image file, e.g. workspace/chart.png or workspace/knowledge_base/diagram.jpg",
                    },
                    "url": {
                        "type": "string",
                        "description": "Alternative: a public URL of an image to upload (e.g. https://example.com/photo.jpg). Use this instead of file_path when the image is not in your workspace.",
                    },
                    "file_name": {
                        "type": "string",
                        "description": "Optional custom filename for the uploaded image. If omitted, the original filename is used.",
                    },
                    "folder": {
                        "type": "string",
                        "description": "Optional CDN folder path, e.g. /agents/reports. Defaults to /clawith.",
                    },
                },
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "discover_resources",
            "description": "Search public MCP registries (Smithery) for tools and capabilities that can extend your abilities. Use this when you encounter a task you cannot handle with your current tools.",
            "parameters": {
                "type": "object",
                "properties": {
                    "query": {
                        "type": "string",
                        "description": "Semantic description of the capability needed, e.g. 'send email', 'query SQL database', 'generate images'",
                    },
                    "max_results": {
                        "type": "integer",
                        "description": "Max results to return (default 5, max 10)",
                    },
                },
                "required": ["query"],
            },
        },
    },
    # ── Feishu Document Tools ──────────────────────────────────────
    {
        "type": "function",
        "function": {
            "name": "feishu_wiki_list",
            "description": (
                "List all sub-pages (child nodes) of a Feishu Wiki (知识库) page. "
                "Works with wiki URLs like 'https://xxx.feishu.cn/wiki/NodeToken'. "
                "Use this when a wiki page has child pages you need to explore. "
                "Returns titles, node_tokens, and obj_tokens for each sub-page. "
                "Each sub-page can then be read with feishu_doc_read using its node_token."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "node_token": {
                        "type": "string",
                        "description": "Wiki node token from the URL, e.g. 'HrGawgXxLiqoS5kT6pUczya3nEc' from 'https://xxx.feishu.cn/wiki/HrGawgXxLiqoS5kT6pUczya3nEc'",
                    },
                    "recursive": {
                        "type": "boolean",
                        "description": "If true, also list sub-pages of sub-pages (up to 2 levels deep). Default false.",
                    },
                },
                "required": ["node_token"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "feishu_doc_read",
            "description": (
                "Read the text content of a Feishu document or Wiki page. "
                "Works with both regular docx URLs (https://xxx.feishu.cn/docx/Token) "
                "and Wiki page URLs (https://xxx.feishu.cn/wiki/Token). "
                "Automatically handles wiki node tokens. "
                "If the page has sub-pages, use feishu_wiki_list to list them."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "document_token": {
                        "type": "string",
                        "description": "Feishu document token (from document URL)",
                    },
                    "max_chars": {
                        "type": "integer",
                        "description": "Max characters to return (default 6000, max 20000)",
                    },
                },
                "required": ["document_token"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "feishu_doc_create",
            "description": "Create a new Feishu document with a given title. Returns the new document token and URL.",
            "parameters": {
                "type": "object",
                "properties": {
                    "title": {
                        "type": "string",
                        "description": "Document title",
                    },
                    "folder_token": {
                        "type": "string",
                        "description": "Optional: parent folder token. Leave empty to create in root My Drive.",
                    },
                },
                "required": ["title"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "feishu_doc_append",
            "description": "Append text content to an existing Feishu document. Content is appended as one or more new paragraphs at the end.",
            "parameters": {
                "type": "object",
                "properties": {
                    "document_token": {
                        "type": "string",
                        "description": "Feishu document token",
                    },
                    "content": {
                        "type": "string",
                        "description": "Text content to append. Supports multiple lines separated by \\n.",
                    },
                },
                "required": ["document_token", "content"],
            },
        },
    },
    # ── Feishu Calendar Tools ──────────────────────────────────────
    {
        "type": "function",
        "function": {
            "name": "feishu_calendar_list",
            "description": "查询飞书日历。**自动读取当前对话用户的真实忙碌时段（freebusy）**，同时列出 bot 创建的日程。用于查询某人是否有空、安排日程时避开冲突。",
            "parameters": {
                "type": "object",
                "properties": {
                    "start_time": {
                        "type": "string",
                        "description": "查询起始时间，ISO 8601 格式，例如 '2026-03-13T00:00:00+08:00'。默认：当前时间。",
                    },
                    "end_time": {
                        "type": "string",
                        "description": "查询截止时间，ISO 8601 格式。默认：7天后。",
                    },
                    "user_open_id": {
                        "type": "string",
                        "description": "要查询 freebusy 的用户 open_id。不填则自动使用当前对话发送者。",
                    },
                    "max_results": {
                        "type": "integer",
                        "description": "Max events to return (default 20)",
                    },
                },
                "required": [],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "feishu_calendar_create",
            "description": "Create a Feishu calendar event immediately. The current user is automatically invited as attendee — no email or authorization required. Just provide the title and time.",
            "parameters": {
                "type": "object",
                "properties": {
                    "summary": {
                        "type": "string",
                        "description": "Event title",
                    },
                    "start_time": {
                        "type": "string",
                        "description": "Event start in ISO 8601 with timezone, e.g. '2026-03-15T14:00:00+08:00'",
                    },
                    "end_time": {
                        "type": "string",
                        "description": "Event end in ISO 8601 with timezone, e.g. '2026-03-15T15:00:00+08:00'",
                    },
                    "description": {
                        "type": "string",
                        "description": "Event description or agenda",
                    },
                    "attendee_names": {
                        "type": "array",
                        "items": {"type": "string"},
                        "description": "Names of colleagues to invite, e.g. ['覃睿', '张三']. Will be looked up automatically via feishu_user_search.",
                    },
                    "attendee_open_ids": {
                        "type": "array",
                        "items": {"type": "string"},
                        "description": "Feishu open_ids to invite directly (if you already have them from feishu_user_search).",
                    },
                    "attendee_emails": {
                        "type": "array",
                        "items": {"type": "string"},
                        "description": "Additional attendee emails to invite (use attendee_names if you only have the name).",
                    },
                    "location": {
                        "type": "string",
                        "description": "Event location or meeting room",
                    },
                    "timezone": {
                        "type": "string",
                        "description": "Timezone, e.g. 'Asia/Shanghai'. Defaults to Asia/Shanghai.",
                    },
                },
                "required": ["summary", "start_time", "end_time"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "feishu_calendar_update",
            "description": "Update an existing Feishu calendar event. Provide only the fields you want to change.",
            "parameters": {
                "type": "object",
                "properties": {
                    "user_email": {"type": "string", "description": "Calendar owner's email"},
                    "event_id": {"type": "string", "description": "Event ID from feishu_calendar_list"},
                    "summary": {"type": "string", "description": "New title"},
                    "description": {"type": "string", "description": "New description"},
                    "start_time": {"type": "string", "description": "New start time (ISO 8601)"},
                    "end_time": {"type": "string", "description": "New end time (ISO 8601)"},
                    "location": {"type": "string", "description": "New location"},
                },
                "required": ["user_email", "event_id"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "feishu_calendar_delete",
            "description": "Delete (cancel) a Feishu calendar event.",
            "parameters": {
                "type": "object",
                "properties": {
                    "user_email": {"type": "string", "description": "Calendar owner's email"},
                    "event_id": {"type": "string", "description": "Event ID to delete"},
                },
                "required": ["user_email", "event_id"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "feishu_doc_share",
            "description": (
                "Manage Feishu document collaborators and permissions. "
                "Can add or remove collaborators with viewer/editor/full_access roles, "
                "or get the current collaborator list. "
                "Accepts colleague names (auto-searched) or open_ids directly."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "document_token": {
                        "type": "string",
                        "description": "Feishu document token (from feishu_doc_create or doc URL)",
                    },
                    "action": {
                        "type": "string",
                        "enum": ["add", "remove", "list"],
                        "description": "'add' to grant access, 'remove' to revoke, 'list' to view current collaborators",
                    },
                    "member_names": {
                        "type": "array",
                        "items": {"type": "string"},
                        "description": "Colleague names to add/remove, e.g. ['覃睿', '张三']. Auto-searched.",
                    },
                    "member_open_ids": {
                        "type": "array",
                        "items": {"type": "string"},
                        "description": "Feishu open_ids to add/remove directly (if already known).",
                    },
                    "permission": {
                        "type": "string",
                        "enum": ["view", "edit", "full_access"],
                        "description": "Permission level: 'view' (read-only), 'edit' (can edit), 'full_access' (can manage). Default: 'edit'",
                    },
                },
                "required": ["document_token", "action"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "feishu_user_search",
            "description": (
                "Search for a colleague in the Feishu (Lark) directory by name. "
                "Returns their open_id, email, and department so you can send messages, "
                "invite them to calendar events, or share documents. "
                "Use this whenever you need to find a colleague's Feishu identity."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "name": {
                        "type": "string",
                        "description": "The colleague's name to search for, e.g. '覃睿' or '张三'",
                    },
                },
                "required": ["name"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "import_mcp_server",
            "description": "Import an MCP server from Smithery registry into the platform. The server's tools become available for use. Use discover_resources first to find the server ID. If previously imported tools stopped working (e.g. OAuth expired), set reauthorize=true to re-run the authorization flow.",
            "parameters": {
                "type": "object",
                "properties": {
                    "server_id": {
                        "type": "string",
                        "description": "Smithery server ID, e.g. '@anthropic/brave-search' or '@anthropic/fetch'",
                    },
                    "config": {
                        "type": "object",
                        "description": "Optional server configuration (e.g. API keys required by the server)",
                    },
                    "reauthorize": {
                        "type": "boolean",
                        "description": "Set to true to force re-authorization of existing tools (e.g. when OAuth token has expired)",
                    },
                },
                "required": ["server_id"],
            },
        },
    },
    # ─── Email Tools ────────────────────────
    {
        "type": "function",
        "function": {
            "name": "send_email",
            "description": "Send an email to one or more recipients. Supports subject, body text, CC, and file attachments from workspace. Requires email configuration in tool settings.",
            "parameters": {
                "type": "object",
                "properties": {
                    "to": {
                        "type": "string",
                        "description": "Recipient email address(es), comma-separated for multiple",
                    },
                    "subject": {
                        "type": "string",
                        "description": "Email subject line",
                    },
                    "body": {
                        "type": "string",
                        "description": "Email body text",
                    },
                    "cc": {
                        "type": "string",
                        "description": "CC recipients, comma-separated (optional)",
                    },
                    "attachments": {
                        "type": "array",
                        "items": {"type": "string"},
                        "description": "List of workspace-relative file paths to attach (optional)",
                    },
                },
                "required": ["to", "subject", "body"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "read_emails",
            "description": "Read emails from your inbox. Can limit the number returned and search by criteria (e.g. FROM, SUBJECT, SINCE date). Requires email configuration in tool settings.",
            "parameters": {
                "type": "object",
                "properties": {
                    "limit": {
                        "type": "integer",
                        "description": "Max number of emails to return (default 10, max 30)",
                    },
                    "search": {
                        "type": "string",
                        "description": "IMAP search criteria, e.g. 'FROM \"john@example.com\"', 'SUBJECT \"meeting\"', 'SINCE 01-Mar-2026'. Default: all emails.",
                    },
                    "folder": {
                        "type": "string",
                        "description": "Mailbox folder, default INBOX",
                    },
                },
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "reply_email",
            "description": "Reply to an email by its Message-ID. Maintains the email thread with proper In-Reply-To headers. Requires email configuration in tool settings.",
            "parameters": {
                "type": "object",
                "properties": {
                    "message_id": {
                        "type": "string",
                        "description": "Message-ID of the email to reply to (from read_emails output)",
                    },
                    "body": {
                        "type": "string",
                        "description": "Reply body text",
                    },
                },
                "required": ["message_id", "body"],
            },
        },
    },
]


# Core tools that should always be available to agents regardless of
# DB configuration.
_ALWAYS_INCLUDE_CORE = {
    "send_channel_file",
    "write_file",
}
# Feishu tools are ONLY included when the agent has a configured Feishu channel,
# to avoid exposing unnecessary tools to non-Feishu agents (reduces hallucination risk).
_FEISHU_TOOL_NAMES = {
    "send_feishu_message",
    "feishu_user_search",
    "feishu_wiki_list",
    "feishu_doc_read",
    "feishu_doc_create",
    "feishu_doc_append",
    "feishu_doc_share",
    "feishu_calendar_list",
    "feishu_calendar_create",
    "feishu_calendar_update",
    "feishu_calendar_delete",
}
_always_core_tools = [t for t in AGENT_TOOLS if t["function"]["name"] in _ALWAYS_INCLUDE_CORE]
_feishu_tools = [t for t in AGENT_TOOLS if t["function"]["name"] in _FEISHU_TOOL_NAMES]


async def _agent_has_feishu(agent_id: uuid.UUID) -> bool:
    """Check if agent has a configured Feishu channel."""
    try:
        from app.models.channel_config import ChannelConfig
        async with async_session() as db:
            r = await db.execute(
                select(ChannelConfig).where(
                    ChannelConfig.agent_id == agent_id,
                    ChannelConfig.channel_type == "feishu",
                    ChannelConfig.is_configured == True,
                )
            )
            return r.scalar_one_or_none() is not None
    except Exception:
        return False


# ─── Dynamic Tool Loading from DB ──────────────────────────────

async def get_agent_tools_for_llm(agent_id: uuid.UUID) -> list[dict]:
    """Load enabled tools for an agent from DB (OpenAI function-calling format).
    
    Falls back to hardcoded AGENT_TOOLS if DB not ready.
    Always includes core system tools (send_channel_file, write_file).
    Feishu tools are only included when the agent has a configured Feishu channel.
    """
    has_feishu = await _agent_has_feishu(agent_id)
    _always_tools = _always_core_tools + (_feishu_tools if has_feishu else [])

    try:
        from app.models.tool import Tool, AgentTool

        async with async_session() as db:
            # Get all globally enabled tools
            all_tools_r = await db.execute(select(Tool).where(Tool.enabled == True))
            all_tools = all_tools_r.scalars().all()

            # Get agent-specific assignments
            agent_tools_r = await db.execute(select(AgentTool).where(AgentTool.agent_id == agent_id))
            assignments = {str(at.tool_id): at for at in agent_tools_r.scalars().all()}

            result = []
            db_tool_names = set()
            for t in all_tools:
                tid = str(t.id)
                at = assignments.get(tid)
                enabled = at.enabled if at else t.is_default
                if not enabled:
                    continue

                # Skip feishu tools if the agent has no Feishu channel configured
                if t.category == "feishu" and not has_feishu:
                    continue

                # Build OpenAI function-calling format
                tool_def = {
                    "type": "function",
                    "function": {
                        "name": t.name,
                        "description": t.description,
                        "parameters": t.parameters_schema or {"type": "object", "properties": {}},
                    },
                }
                result.append(tool_def)
                db_tool_names.add(t.name)

            if result:
                # Append always-available system tools that aren't already in the DB list
                for t in _always_tools:
                    if t["function"]["name"] not in db_tool_names:
                        result.append(t)
                return result
    except Exception as e:
        print(f"[Tools] DB load failed, using fallback: {e}")

    # Fallback to hardcoded tools
    return AGENT_TOOLS


# ─── Workspace initialization ──────────────────────────────────

async def ensure_workspace(agent_id: uuid.UUID) -> Path:
    """Initialize agent workspace with standard structure."""
    ws = WORKSPACE_ROOT / str(agent_id)
    ws.mkdir(parents=True, exist_ok=True)

    # Create standard directories
    (ws / "skills").mkdir(exist_ok=True)
    (ws / "workspace").mkdir(exist_ok=True)
    (ws / "workspace" / "knowledge_base").mkdir(exist_ok=True)
    (ws / "memory").mkdir(exist_ok=True)

    # Ensure shared enterprise_info directory exists
    enterprise_dir = WORKSPACE_ROOT / "enterprise_info"
    enterprise_dir.mkdir(parents=True, exist_ok=True)
    (enterprise_dir / "knowledge_base").mkdir(exist_ok=True)
    # Create default company profile if missing
    profile_path = enterprise_dir / "company_profile.md"
    if not profile_path.exists():
        profile_path.write_text("# Company Profile\n\n_Edit company information here. All digital employees can access this._\n\n## Basic Info\n- Company Name:\n- Industry:\n- Founded:\n\n## Business Overview\n\n## Organization Structure\n\n## Company Culture\n", encoding="utf-8")

    # Migrate: move root-level memory.md into memory/ directory
    if (ws / "memory.md").exists() and not (ws / "memory" / "memory.md").exists():
        import shutil
        shutil.move(str(ws / "memory.md"), str(ws / "memory" / "memory.md"))

    # Create default memory file if missing
    if not (ws / "memory" / "memory.md").exists():
        (ws / "memory" / "memory.md").write_text("# Memory\n\n_Record important information and knowledge here._\n", encoding="utf-8")

    if not (ws / "soul.md").exists():
        # Try to load from DB
        try:
            from app.models.agent import Agent
            async with async_session() as db:
                r = await db.execute(select(Agent).where(Agent.id == agent_id))
                agent = r.scalar_one_or_none()
                if agent and agent.role_description:
                    (ws / "soul.md").write_text(
                        f"# Personality\n\n{agent.role_description}\n",
                        encoding="utf-8",
                    )
                else:
                    (ws / "soul.md").write_text("# Personality\n\n_Describe your role and responsibilities._\n", encoding="utf-8")
        except Exception:
            (ws / "soul.md").write_text("# Personality\n\n_Describe your role and responsibilities._\n", encoding="utf-8")

    # Always sync tasks from DB
    await _sync_tasks_to_file(agent_id, ws)

    return ws


async def _sync_tasks_to_file(agent_id: uuid.UUID, ws: Path):
    """Sync tasks from DB to tasks.json in workspace."""
    try:
        async with async_session() as db:
            result = await db.execute(
                select(Task).where(Task.agent_id == agent_id).order_by(Task.created_at.desc())
            )
            tasks = result.scalars().all()

        task_list = []
        for t in tasks:
            task_list.append({
                "title": t.title,
                "status": t.status,
                "priority": t.priority,
                "description": t.description or "",
                "created_at": t.created_at.isoformat() if t.created_at else "",
                "completed_at": t.completed_at.isoformat() if t.completed_at else "",
            })

        (ws / "tasks.json").write_text(
            json.dumps(task_list, ensure_ascii=False, indent=2),
            encoding="utf-8",
        )
    except Exception as e:
        print(f"[AgentTools] Failed to sync tasks: {e}")


# ─── Tool Executors ─────────────────────────────────────────────

# Mapping from tool_name to autonomy action_type
_TOOL_AUTONOMY_MAP = {
    "write_file": "write_workspace_files",
    "delete_file": "delete_files",
    "send_feishu_message": "send_feishu_message",
    "send_message_to_agent": "send_feishu_message",
    "web_search": "web_search",
    "execute_code": "execute_code",
}


async def _execute_tool_direct(
    tool_name: str,
    arguments: dict,
    agent_id: uuid.UUID,
) -> str:
    """Execute a tool directly, bypassing autonomy checks.

    Used by the approval post-processing hook after an action
    has been approved and needs to actually run.
    """
    ws = await ensure_workspace(agent_id)
    try:
        if tool_name == "delete_file":
            return _delete_file(ws, arguments.get("path", ""))
        elif tool_name == "write_file":
            path = arguments.get("path")
            content = arguments.get("content", "")
            if not path:
                return "Missing path"
            return _write_file(ws, path, content)
        elif tool_name == "execute_code":
            return await _execute_code(agent_id, ws, arguments)
        elif tool_name == "web_search":
            return await _web_search(arguments)
        elif tool_name == "jina_search":
            return await _jina_search(arguments)
        elif tool_name == "send_feishu_message":
            return await _send_feishu_message(agent_id, arguments)
        elif tool_name == "send_message_to_agent":
            return await _send_message_to_agent(agent_id, arguments)
        else:
            return f"Tool {tool_name} does not support post-approval execution"
    except Exception as e:
        return f"Error executing {tool_name}: {e}"


async def execute_tool(
    tool_name: str,
    arguments: dict,
    agent_id: uuid.UUID,
    user_id: uuid.UUID,
) -> str:
    """Execute a tool call and return the result as a string."""
    ws = await ensure_workspace(agent_id)

    # ── Autonomy boundary check ──
    action_type = _TOOL_AUTONOMY_MAP.get(tool_name)
    if action_type:
        try:
            from app.services.autonomy_service import autonomy_service
            from app.models.agent import Agent as AgentModel
            async with async_session() as _adb:
                _ar = await _adb.execute(select(AgentModel).where(AgentModel.id == agent_id))
                _agent = _ar.scalar_one_or_none()
                if _agent:
                    result_check = await autonomy_service.check_and_enforce(
                        _adb, _agent, action_type, {"tool": tool_name, "args": str(arguments)[:200], "requested_by": str(user_id)}
                    )
                    await _adb.commit()
                    if not result_check.get("allowed"):
                        level = result_check.get("level", "L3")
                        if level == "L3":
                            return f"⏳ This action requires approval. An approval request has been sent. Please wait for approval before retrying. (Approval ID: {result_check.get('approval_id', 'N/A')})"
                        return f"❌ Action denied: {result_check.get('message', 'unknown reason')}"
        except Exception as e:
            print(f"[Autonomy] Check failed — blocking as safety measure: {e}")
            return f"⚠️ Autonomy check failed ({e}). Operation blocked for safety. Please retry or contact admin."

    try:
        if tool_name == "list_files":
            result = _list_files(ws, arguments.get("path", ""))
        elif tool_name == "read_file":
            path = arguments.get("path")
            if not path:
                return "❌ Missing required argument 'path' for read_file"
            result = _read_file(ws, path)
        elif tool_name == "read_document":
            path = arguments.get("path")
            if not path:
                return "❌ Missing required argument 'path' for read_document"
            max_chars = min(int(arguments.get("max_chars", 8000)), 20000)
            result = await _read_document(ws, path, max_chars=max_chars)
        elif tool_name == "write_file":
            path = arguments.get("path")
            content = arguments.get("content")
            if not path:
                return "❌ Missing required argument 'path' for write_file. Please provide a file path like 'skills/my-skill/SKILL.md'"
            if content is None:
                return "❌ Missing required argument 'content' for write_file"
            result = _write_file(ws, path, content)
        elif tool_name == "delete_file":
            result = _delete_file(ws, arguments.get("path", ""))
        elif tool_name == "manage_tasks":
            result = await _manage_tasks(agent_id, user_id, ws, arguments)
        elif tool_name == "set_trigger":
            result = await _handle_set_trigger(agent_id, arguments)
        elif tool_name == "update_trigger":
            result = await _handle_update_trigger(agent_id, arguments)
        elif tool_name == "cancel_trigger":
            result = await _handle_cancel_trigger(agent_id, arguments)
        elif tool_name == "list_triggers":
            result = await _handle_list_triggers(agent_id)
        elif tool_name == "send_feishu_message":
            result = await _send_feishu_message(agent_id, arguments)
        elif tool_name == "send_web_message":
            result = await _send_web_message(agent_id, arguments)
        elif tool_name == "send_message_to_agent":
            result = await _send_message_to_agent(agent_id, arguments)
        elif tool_name == "send_channel_file":
            result = await _send_channel_file(agent_id, ws, arguments)
        elif tool_name == "web_search":
            result = await _web_search(arguments)
        elif tool_name == "jina_search":
            result = await _jina_search(arguments)
        elif tool_name == "bing_search":
            result = await _jina_search(arguments)  # redirect legacy to jina
        elif tool_name == "jina_read":
            result = await _jina_read(arguments)
        elif tool_name == "read_webpage":
            result = await _jina_read(arguments)  # redirect legacy to jina
        elif tool_name == "plaza_get_new_posts":
            result = await _plaza_get_new_posts(arguments)
        elif tool_name == "plaza_create_post":
            result = await _plaza_create_post(agent_id, arguments)
        elif tool_name == "plaza_add_comment":
            result = await _plaza_add_comment(agent_id, arguments)
        elif tool_name == "execute_code":
            result = await _execute_code(ws, arguments)
        elif tool_name == "upload_image":
            result = await _upload_image(agent_id, ws, arguments)
        elif tool_name == "discover_resources":
            result = await _discover_resources(arguments)
        elif tool_name == "import_mcp_server":
            result = await _import_mcp_server(agent_id, arguments)
        # ── Feishu Document Tools ──
        elif tool_name == "feishu_wiki_list":
            result = await _feishu_wiki_list(agent_id, arguments)
        elif tool_name == "feishu_doc_read":
            result = await _feishu_doc_read(agent_id, arguments)
        elif tool_name == "feishu_doc_create":
            result = await _feishu_doc_create(agent_id, arguments)
        elif tool_name == "feishu_doc_append":
            result = await _feishu_doc_append(agent_id, arguments)
        # ── Feishu Calendar Tools ──
        elif tool_name == "feishu_doc_share":
            result = await _feishu_doc_share(agent_id, arguments)
        elif tool_name == "feishu_user_search":
            result = await _feishu_user_search(agent_id, arguments)
        elif tool_name == "feishu_calendar_list":
            result = await _feishu_calendar_list(agent_id, arguments)
        elif tool_name == "feishu_calendar_create":
            result = await _feishu_calendar_create(agent_id, arguments)
        elif tool_name == "feishu_calendar_update":
            result = await _feishu_calendar_update(agent_id, arguments)
        elif tool_name == "feishu_calendar_delete":
            result = await _feishu_calendar_delete(agent_id, arguments)
        # ── Email Tools ──
        elif tool_name in ("send_email", "read_emails", "reply_email"):
            result = await _handle_email_tool(tool_name, agent_id, ws, arguments)
        else:
            # Try MCP tool execution
            result = await _execute_mcp_tool(tool_name, arguments, agent_id=agent_id)

        # Log tool call activity (skip noisy read operations)
        if tool_name not in ("list_files", "read_file", "read_document"):
            from app.services.activity_logger import log_activity
            await log_activity(
                agent_id, "tool_call",
                f"Called tool {tool_name}: {result[:80]}",
                detail={"tool": tool_name, "args": {k: str(v)[:100] for k, v in arguments.items()}, "result": result[:300]},
            )
        return result
    except Exception as e:
        import traceback
        traceback.print_exc()
        return f"Tool execution error ({tool_name}): {type(e).__name__}: {str(e)[:200]}"


async def _web_search(arguments: dict) -> str:
    """Search the web using a configurable search engine (reads config from DB)."""
    import httpx
    import re

    query = arguments.get("query", "")
    if not query:
        return "❌ Please provide search keywords"

    # Load config from DB
    config = {}
    try:
        from app.models.tool import Tool
        async with async_session() as db:
            r = await db.execute(select(Tool).where(Tool.name == "web_search"))
            tool = r.scalar_one_or_none()
            if tool and tool.config:
                config = tool.config
    except Exception:
        pass

    engine = config.get("search_engine", "duckduckgo")
    api_key = config.get("api_key", "")
    max_results = min(arguments.get("max_results", config.get("max_results", 5)), 10)
    language = config.get("language", "zh-CN")

    try:
        if engine == "tavily" and api_key:
            return await _search_tavily(query, api_key, max_results)
        elif engine == "google" and api_key:
            return await _search_google(query, api_key, max_results, language)
        elif engine == "bing" and api_key:
            return await _search_bing(query, api_key, max_results, language)
        else:
            return await _search_duckduckgo(query, max_results)
    except Exception as e:
        return f"❌ Search error ({engine}): {str(e)[:200]}"


async def _search_duckduckgo(query: str, max_results: int) -> str:
    """Search via DuckDuckGo HTML (free, no API key)."""
    import httpx, re

    async with httpx.AsyncClient(follow_redirects=True) as client:
        resp = await client.get(
            "https://html.duckduckgo.com/html/",
            params={"q": query},
            headers={"User-Agent": "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7)"},
            timeout=10,
        )

    results = []
    blocks = re.findall(
        r'<a[^>]*class="result__a"[^>]*href="([^"]*)"[^>]*>(.*?)</a>.*?'
        r'<a[^>]*class="result__snippet"[^>]*>(.*?)</a>',
        resp.text, re.DOTALL,
    )
    for url, title, snippet in blocks[:max_results]:
        title = re.sub(r'<[^>]+>', '', title).strip()
        snippet = re.sub(r'<[^>]+>', '', snippet).strip()
        if "uddg=" in url:
            from urllib.parse import unquote, parse_qs, urlparse
            parsed = parse_qs(urlparse(url).query)
            url = unquote(parsed.get("uddg", [url])[0])
        results.append(f"**{title}**\n{url}\n{snippet}")

    if not results:
        return f'🔍 No results found for "{query}"'
    return f'🔍 DuckDuckGo results for "{query}" ({len(results)} items):\n\n' + "\n\n---\n\n".join(results)

async def _get_jina_api_key() -> str:
    """Read Jina API key from DB system_settings first, then fall back to env."""
    try:
        from app.database import async_session
        from app.models.system_settings import SystemSetting
        from sqlalchemy import select
        async with async_session() as db:
            result = await db.execute(select(SystemSetting).where(SystemSetting.key == "jina_api_key"))
            setting = result.scalar_one_or_none()
            if setting and setting.value.get("api_key"):
                return setting.value["api_key"]
    except Exception:
        pass
    from app.config import get_settings
    return get_settings().JINA_API_KEY


async def _jina_search(arguments: dict) -> str:
    """Search via Jina AI Search API (s.jina.ai). Returns full content per result, not just snippets."""
    import httpx

    query = arguments.get("query", "").strip()
    if not query:
        return "❌ Please provide search keywords"

    max_results = min(arguments.get("max_results", 5), 10)
    api_key = await _get_jina_api_key()

    headers: dict = {
        "Accept": "application/json",
        "X-Respond-With": "no-content",  # return snippets/descriptions, not full pages (faster)
        "X-Return-Format": "markdown",
    }
    if api_key:
        headers["Authorization"] = f"Bearer {api_key}"

    try:
        async with httpx.AsyncClient(follow_redirects=True, timeout=30) as client:
            resp = await client.get(
                f"https://s.jina.ai/{__import__('urllib.parse', fromlist=['quote']).quote(query)}",
                headers=headers,
            )

        if resp.status_code != 200:
            return f"❌ Jina Search error HTTP {resp.status_code}: {resp.text[:200]}"

        data = resp.json()
        items = data.get("data", [])[:max_results]

        if not items:
            return f'🔍 No results found for "{query}"'

        parts = []
        for i, item in enumerate(items, 1):
            title = item.get("title", "Untitled")
            url = item.get("url", "")
            description = item.get("description", "") or item.get("content", "")[:500]
            parts.append(f"**{i}. {title}**\n{url}\n{description}")

        return f'🔍 Jina Search results for "{query}" ({len(items)} items):\n\n' + "\n\n---\n\n".join(parts)

    except Exception as e:
        return f"❌ Jina Search error: {str(e)[:300]}"


async def _jina_read(arguments: dict) -> str:
    """Read web page via Jina AI Reader API (r.jina.ai). Returns clean structured markdown."""
    import httpx
    from app.config import get_settings

    url = arguments.get("url", "").strip()
    if not url:
        return "❌ Please provide a URL"
    if not url.startswith("http"):
        url = "https://" + url

    max_chars = min(arguments.get("max_chars", 8000), 20000)
    api_key = await _get_jina_api_key()

    headers: dict = {
        "Accept": "text/plain, text/markdown, */*",
        "X-Return-Format": "markdown",
        "X-Remove-Selector": "header, footer, nav, aside, .ads, .advertisement",
    }
    if api_key:
        headers["Authorization"] = f"Bearer {api_key}"

    try:
        async with httpx.AsyncClient(follow_redirects=True, timeout=30) as client:
            resp = await client.get(
                f"https://r.jina.ai/{url}",
                headers=headers,
            )

        if resp.status_code != 200:
            return f"❌ Jina Reader error HTTP {resp.status_code}: {resp.text[:200]}"

        text = resp.text.strip()
        if not text or len(text) < 100:
            return f"❌ Jina Reader returned empty content for {url}"

        if len(text) > max_chars:
            text = text[:max_chars] + f"\n\n[... truncated at {max_chars} chars]"

        return f"📄 **Content from: {url}**\n\n{text}"

    except Exception as e:
        return f"❌ Jina Reader error: {str(e)[:300]}"



async def _search_tavily(query: str, api_key: str, max_results: int) -> str:
    """Search via Tavily API (AI-optimized search)."""
    import httpx

    async with httpx.AsyncClient() as client:
        resp = await client.post(
            "https://api.tavily.com/search",
            json={"query": query, "max_results": max_results, "search_depth": "basic"},
            headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"},
            timeout=15,
        )
        data = resp.json()

    if "results" not in data:
        return f"❌ Tavily search failed: {data.get('error', str(data)[:200])}"

    results = []
    for r in data["results"][:max_results]:
        results.append(f"**{r.get('title', '')}**\n{r.get('url', '')}\n{r.get('content', '')[:200]}")

    if not results:
        return f'🔍 No results found for "{query}"'
    return f'🔍 Tavily search for "{query}" ({len(results)} items):\n\n' + "\n\n---\n\n".join(results)


async def _search_google(query: str, api_key: str, max_results: int, language: str) -> str:
    """Search via Google Custom Search JSON API."""
    import httpx

    # api_key format: "API_KEY:CX_ID"
    parts = api_key.split(":", 1)
    if len(parts) != 2:
        return "❌ Google search requires API key in format 'API_KEY:SEARCH_ENGINE_ID'"

    gapi_key, cx = parts
    async with httpx.AsyncClient() as client:
        resp = await client.get(
            "https://www.googleapis.com/customsearch/v1",
            params={"key": gapi_key, "cx": cx, "q": query, "num": max_results, "lr": f"lang_{language[:2]}"},
            timeout=10,
        )
        data = resp.json()

    results = []
    for item in data.get("items", [])[:max_results]:
        results.append(f"**{item.get('title', '')}**\n{item.get('link', '')}\n{item.get('snippet', '')}")

    if not results:
        return f'🔍 No results found for "{query}"'
    return f'🔍 Google search for "{query}" ({len(results)} items):\n\n' + "\n\n---\n\n".join(results)


async def _search_bing(query: str, api_key: str, max_results: int, language: str) -> str:
    """Search via Bing Web Search API."""
    import httpx

    async with httpx.AsyncClient() as client:
        resp = await client.get(
            "https://api.bing.microsoft.com/v7.0/search",
            params={"q": query, "count": max_results, "mkt": language},
            headers={"Ocp-Apim-Subscription-Key": api_key},
            timeout=10,
        )
        data = resp.json()

    results = []
    for item in data.get("webPages", {}).get("value", [])[:max_results]:
        results.append(f"**{item.get('name', '')}**\n{item.get('url', '')}\n{item.get('snippet', '')}")

    if not results:
        return f'🔍 No results found for "{query}"'
    return f'🔍 Bing search for "{query}" ({len(results)} items):\n\n' + "\n\n---\n\n".join(results)


async def _send_channel_file(agent_id: uuid.UUID, ws: Path, arguments: dict) -> str:
    """Send a file to the user via the current channel or return a download URL for web chat."""
    rel_path = arguments.get("file_path", "").strip()
    accompany_msg = arguments.get("message", "")
    if not rel_path:
        return "❌ file_path is required"

    # Resolve file path within agent workspace
    file_path = (ws / rel_path).resolve()
    ws_resolved = ws.resolve()
    if not str(file_path).startswith(str(ws_resolved)):
        # Also allow workspace/ prefix pointing to same location
        file_path = (WORKSPACE_ROOT / str(agent_id) / rel_path).resolve()
        if not file_path.exists():
            return f"❌ File not found: {rel_path}"
    if not file_path.exists():
        return f"❌ File not found: {rel_path}"

    sender = channel_file_sender.get()
    if sender is not None:
        # Channel mode: call the channel-specific send function
        try:
            await sender(file_path, accompany_msg)
            return f"✅ File '{file_path.name}' sent to user via channel."
        except Exception as e:
            return f"❌ Failed to send file: {e}"
    else:
        # Web chat mode: return a download URL
        aid = channel_web_agent_id.get() or str(agent_id)
        base_abs = (WORKSPACE_ROOT / str(agent_id)).resolve()
        try:
            file_rel = str(file_path.resolve().relative_to(base_abs))
        except ValueError:
            file_rel = rel_path
        from app.config import get_settings as _gs
        _s = _gs()
        base_url = getattr(_s, 'BASE_URL', '').rstrip('/') or ''
        download_url = f"{base_url}/api/agents/{aid}/files/download?path={file_rel}"
        msg = f"✅ File ready: [{file_path.name}]({download_url})"
        if accompany_msg:
            msg = accompany_msg + "\n\n" + msg
        return msg


async def _execute_mcp_tool(tool_name: str, arguments: dict, agent_id=None) -> str:
    """Execute a tool via MCP if it exists in the DB as an MCP tool."""
    try:
        from app.models.tool import Tool, AgentTool
        from app.services.mcp_client import MCPClient

        async with async_session() as db:
            result = await db.execute(select(Tool).where(Tool.name == tool_name, Tool.type == "mcp"))
            tool = result.scalar_one_or_none()
            # Load per-agent config override
            agent_config = {}
            if tool and agent_id:
                at_r = await db.execute(
                    select(AgentTool).where(
                        AgentTool.agent_id == agent_id,
                        AgentTool.tool_id == tool.id,
                    )
                )
                at = at_r.scalar_one_or_none()
                agent_config = (at.config or {}) if at else {}

        if not tool:
            return f"Unknown tool: {tool_name}"

        if not tool.mcp_server_url:
            return f"❌ MCP tool {tool_name} has no server URL configured"

        # Merge global config + agent override
        merged_config = {**(tool.config or {}), **agent_config}

        mcp_url = tool.mcp_server_url
        mcp_name = tool.mcp_tool_name or tool_name

        # Detect Smithery-hosted MCP servers (*.run.tools URLs)
        # These need Smithery Connect to route tool calls
        if ".run.tools" in mcp_url and merged_config:
            return await _execute_via_smithery_connect(mcp_url, mcp_name, arguments, merged_config, agent_id=agent_id)

        # Direct MCP call for non-Smithery servers
        # Priority for API key:
        # 1. Per-agent tool config (api_key / atlassian_api_key)
        # 2. Agent's Atlassian channel config (for atlassian_* tools)
        direct_api_key = merged_config.get("api_key") or merged_config.get("atlassian_api_key")
        if not direct_api_key and tool.mcp_server_name == "Atlassian Rovo":
            try:
                from app.api.atlassian import get_atlassian_api_key_for_agent
                direct_api_key = await get_atlassian_api_key_for_agent(agent_id)
            except Exception:
                pass
        client = MCPClient(mcp_url, api_key=direct_api_key)
        return await client.call_tool(mcp_name, arguments)

    except Exception as e:
        return f"❌ MCP tool execution error: {str(e)[:200]}"


async def _execute_via_smithery_connect(mcp_url: str, tool_name: str, arguments: dict, config: dict, agent_id=None) -> str:
    """Execute an MCP tool via Smithery Connect API.

    Uses stored namespace/connection or falls back to creating one.
    Smithery Connect returns SSE-format responses that need special parsing.
    """
    import httpx
    import json as json_mod

    # Get Smithery API key centrally (from discover_resources/import_mcp_server AgentTool config)
    from app.services.resource_discovery import _get_smithery_api_key
    api_key = await _get_smithery_api_key(agent_id)
    if not api_key:
        return (
            "❌ Smithery API key not configured.\n\n"
            "请提供你的 Smithery API Key，你可以通过以下步骤获取：\n"
            "1. 注册/登录 https://smithery.ai\n"
            "2. 前往 https://smithery.ai/account/api-keys 创建 API Key\n"
            "3. 将 Key 提供给我，我会帮你配置"
        )

    # Get namespace + connection from tool config, or use defaults
    namespace = config.pop("smithery_namespace", None)
    connection_id = config.pop("smithery_connection_id", None)

    if not namespace or not connection_id:
        # Fallback: try to get from Smithery settings
        try:
            from app.models.tool import Tool
            async with async_session() as db:
                r = await db.execute(select(Tool).where(Tool.name == "discover_resources"))
                disc_tool = r.scalar_one_or_none()
                if disc_tool and disc_tool.config:
                    namespace = namespace or disc_tool.config.get("smithery_namespace")
                    connection_id = connection_id or disc_tool.config.get("smithery_connection_id")
        except Exception:
            pass

    if not namespace or not connection_id:
        return (
            "❌ Smithery Connect namespace/connection not configured. "
            "Please set smithery_namespace and smithery_connection_id in the tool configuration."
        )

    headers = {
        "Authorization": f"Bearer {api_key}",
        "Content-Type": "application/json",
    }

    try:
        async with httpx.AsyncClient(timeout=30, follow_redirects=True) as client:
            # Call the tool via the existing connection
            tool_resp = await client.post(
                f"https://api.smithery.ai/connect/{namespace}/{connection_id}/mcp",
                json={
                    "jsonrpc": "2.0",
                    "id": 1,
                    "method": "tools/call",
                    "params": {
                        "name": tool_name,
                        "arguments": arguments,
                    },
                },
                headers=headers,
            )

            # Detect auth/connection failures and attempt auto-recovery
            if tool_resp.status_code in (401, 403, 404):
                recovery_result = await _smithery_auto_recover(
                    api_key, mcp_url, namespace, connection_id, agent_id
                )
                if recovery_result:
                    return recovery_result
                # If recovery returned None, fall through to normal parsing

            # Smithery Connect returns SSE format: "event: message\ndata: {...}\n"
            raw = tool_resp.text
            data = None

            # Parse SSE response
            for line in raw.split("\n"):
                line = line.strip()
                if line.startswith("data: "):
                    try:
                        data = json_mod.loads(line[6:])
                        break
                    except json_mod.JSONDecodeError:
                        pass

            # Fallback: try parsing as plain JSON
            if data is None:
                try:
                    data = json_mod.loads(raw)
                except json_mod.JSONDecodeError:
                    return f"❌ Unexpected response from Smithery: {raw[:300]}"

            if "error" in data:
                err = data["error"]
                msg = err.get("message", str(err)) if isinstance(err, dict) else str(err)
                # Check if error indicates auth/connection issue
                auth_keywords = ["auth", "unauthorized", "forbidden", "expired", "not found", "connection"]
                if any(kw in msg.lower() for kw in auth_keywords):
                    recovery_result = await _smithery_auto_recover(
                        api_key, mcp_url, namespace, connection_id, agent_id
                    )
                    if recovery_result:
                        return recovery_result
                return f"❌ MCP tool error: {msg[:300]}"

            result = data.get("result", {})
            if isinstance(result, str):
                return result

            content_blocks = result.get("content", []) if isinstance(result, dict) else []
            texts = []
            for block in content_blocks:
                if isinstance(block, str):
                    texts.append(block)
                elif isinstance(block, dict):
                    if block.get("type") == "text":
                        texts.append(block.get("text", ""))
                    elif block.get("type") == "image":
                        texts.append(f"[Image: {block.get('mimeType', 'image')}]")
                    else:
                        texts.append(str(block))
                else:
                    texts.append(str(block))

            return "\n".join(texts) if texts else str(result)

    except Exception as e:
        return f"❌ Smithery Connect error: {str(e)[:200]}"


async def _smithery_auto_recover(api_key: str, mcp_url: str, namespace: str, connection_id: str, agent_id=None) -> str | None:
    """Attempt to auto-recover a failed Smithery connection.

    Re-creates the Smithery Connect connection. If OAuth is needed,
    returns the auth URL for the user. Returns None if recovery fails silently.
    """
    try:
        from app.services.resource_discovery import _ensure_smithery_connection
        display_name = connection_id.replace("-", " ").title() if connection_id else "MCP Server"

        conn_result = await _ensure_smithery_connection(api_key, mcp_url, display_name)
        if "error" in conn_result:
            return (
                f"❌ MCP tool connection expired and auto-recovery failed: {conn_result['error']}\n\n"
                f"💡 Please re-authorize by telling me: `import_mcp_server(server_id=\"...\", reauthorize=true)`"
            )

        # Update stored config with new connection info
        new_config = {
            "smithery_namespace": conn_result["namespace"],
            "smithery_connection_id": conn_result["connection_id"],
        }
        if agent_id:
            try:
                from app.models.tool import Tool, AgentTool
                async with async_session() as db:
                    # Update all MCP tools for this server URL
                    r = await db.execute(
                        select(Tool).where(Tool.mcp_server_url == mcp_url, Tool.type == "mcp")
                    )
                    for tool in r.scalars().all():
                        at_r = await db.execute(
                            select(AgentTool).where(
                                AgentTool.agent_id == agent_id,
                                AgentTool.tool_id == tool.id,
                            )
                        )
                        at = at_r.scalar_one_or_none()
                        if at:
                            at.config = {**(at.config or {}), **new_config}
                    await db.commit()
            except Exception:
                pass  # Non-critical — connection may still work

        if conn_result.get("auth_url"):
            return (
                f"🔐 MCP tool connection expired. Re-authorization needed.\n\n"
                f"Please visit the following URL to re-authorize:\n"
                f"{conn_result['auth_url']}\n\n"
                f"After completing authorization, the tools will work again automatically."
            )

        # Connection re-created without OAuth — should work now
        return None  # Signal caller to retry (but we don't retry here to avoid loops)

    except Exception as e:
        return f"❌ Auto-recovery failed: {str(e)[:200]}"


def _list_files(ws: Path, rel_path: str) -> str:
    # Handle enterprise_info/ as shared directory
    if rel_path and rel_path.startswith("enterprise_info"):
        target = (WORKSPACE_ROOT / rel_path).resolve()
        enterprise_root = (WORKSPACE_ROOT / "enterprise_info").resolve()
        if not str(target).startswith(str(enterprise_root)):
            return "Access denied for this path"
    else:
        target = (ws / rel_path) if rel_path else ws
        target = target.resolve()
        if not str(target).startswith(str(ws.resolve())):
            return "Access denied for this path"

    if not target.exists():
        return f"Directory not found: {rel_path or '/'}"

    items = []
    # If listing root, also show enterprise_info entry
    if not rel_path:
        enterprise_dir = WORKSPACE_ROOT / "enterprise_info"
        if enterprise_dir.exists():
            items.append("  📁 enterprise_info/ (shared company info)")

    dir_count = 0
    file_count = 0
    for p in sorted(target.iterdir()):
        if p.name.startswith("."):
            continue
        if p.is_dir():
            dir_count += 1
            child_count = len([c for c in p.iterdir() if not c.name.startswith(".")])
            items.append(f"  📁 {p.name}/ ({child_count} items)")
        elif p.is_file():
            file_count += 1
            size_bytes = p.stat().st_size
            if size_bytes < 1024:
                size_str = f"{size_bytes}B"
            else:
                size_str = f"{size_bytes/1024:.1f}KB"
            items.append(f"  📄 {p.name} ({size_str})")

    if not items:
        return f"📂 {rel_path or 'root'}: Empty directory (0 files, 0 folders)"

    header = f"📂 {rel_path or 'root'}: {dir_count} folder(s), {file_count} file(s)\n"
    return header + "\n".join(items)


def _read_file(ws: Path, rel_path: str) -> str:
    # Handle enterprise_info/ as shared directory
    if rel_path and rel_path.startswith("enterprise_info"):
        file_path = (WORKSPACE_ROOT / rel_path).resolve()
        enterprise_root = (WORKSPACE_ROOT / "enterprise_info").resolve()
        if not str(file_path).startswith(str(enterprise_root)):
            return "Access denied for this path"
    else:
        file_path = (ws / rel_path).resolve()
        if not str(file_path).startswith(str(ws.resolve())):
            return "Access denied for this path"

    if not file_path.exists():
        return f"File not found: {rel_path}"

    try:
        content = file_path.read_text(encoding="utf-8", errors="replace")
        if len(content) > 6000:
            content = content[:6000] + f"\n\n...[truncated, {len(content)} chars total]"
        return content
    except Exception as e:
        return f"Read failed: {e}"


async def _read_document(ws: Path, rel_path: str, max_chars: int = 8000) -> str:
    """Read content from office documents (PDF, DOCX, XLSX, PPTX)."""
    # Handle enterprise_info/ as shared directory
    if rel_path and rel_path.startswith("enterprise_info"):
        file_path = (WORKSPACE_ROOT / rel_path).resolve()
        enterprise_root = (WORKSPACE_ROOT / "enterprise_info").resolve()
        if not str(file_path).startswith(str(enterprise_root)):
            return "Access denied for this path"
    else:
        file_path = (ws / rel_path).resolve()
        if not str(file_path).startswith(str(ws.resolve())):
            return "Access denied for this path"

    if not file_path.exists():
        return f"File not found: {rel_path}"

    ext = file_path.suffix.lower()
    try:
        if ext == ".pdf":
            import pdfplumber
            text_parts = []
            with pdfplumber.open(str(file_path)) as pdf:
                for i, page in enumerate(pdf.pages[:50]):  # Limit to 50 pages
                    page_text = page.extract_text() or ""
                    if page_text:
                        text_parts.append(f"--- Page {i+1} ---\n{page_text}")
            content = "\n\n".join(text_parts) if text_parts else "(PDF is empty or text extraction failed)"

        elif ext == ".docx":
            from docx import Document
            from docx.oxml.ns import qn
            doc = Document(str(file_path))
            lines: list[str] = []

            def _extract_para_text(para) -> str:
                return para.text.strip()

            def _extract_table(table) -> str:
                """Flatten a table into readable text."""
                rows = []
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    # Remove duplicate adjacent cells (merged cells repeat)
                    deduped = [cells[0]] + [c for i, c in enumerate(cells[1:]) if c != cells[i]]
                    row_str = " | ".join(c for c in deduped if c)
                    if row_str:
                        rows.append(row_str)
                return "\n".join(rows)

            # 1. Main paragraphs
            for para in doc.paragraphs:
                t = _extract_para_text(para)
                if t:
                    lines.append(t)

            # 2. Tables in main body
            for table in doc.tables:
                t = _extract_table(table)
                if t:
                    lines.append(t)

            # 3. Text boxes / drawing shapes (wmf/shapes in body XML)
            for shape in doc.element.body.iter(qn("w:txbxContent")):
                for child in shape.iter(qn("w:t")):
                    if child.text and child.text.strip():
                        lines.append(child.text.strip())

            # 4. Headers and footers
            for section in doc.sections:
                for hf in [section.header, section.footer]:
                    if hf and hf.is_linked_to_previous is False:
                        for para in hf.paragraphs:
                            t = para.text.strip()
                            if t:
                                lines.append(t)

            content = "\n".join(lines) if lines else "(Document is empty or uses unsupported formatting)"

        elif ext == ".xlsx":
            from openpyxl import load_workbook
            wb = load_workbook(str(file_path), read_only=True, data_only=True)
            sheets = []
            for ws_name in wb.sheetnames[:10]:  # Limit to 10 sheets
                sheet = wb[ws_name]
                rows = []
                for row in sheet.iter_rows(max_row=200, values_only=True):
                    row_str = "\t".join(str(c) if c is not None else "" for c in row)
                    if row_str.strip():
                        rows.append(row_str)
                if rows:
                    sheets.append(f"=== Sheet: {ws_name} ===\n" + "\n".join(rows))
            wb.close()
            content = "\n\n".join(sheets) if sheets else "(Excel is empty)"

        elif ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(str(file_path))
            slides = []
            for i, slide in enumerate(prs.slides[:50]):
                texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text)
                if texts:
                    slides.append(f"--- Slide {i+1} ---\n" + "\n".join(texts))
            content = "\n\n".join(slides) if slides else "(PPT is empty)"

        elif ext in (".txt", ".md", ".json", ".csv", ".log"):
            content = file_path.read_text(encoding="utf-8", errors="replace")

        else:
            return f"Unsupported file format: {ext}. Supported: PDF, DOCX, XLSX, PPTX, TXT, MD, CSV"

        if len(content) > max_chars:
            content = content[:max_chars] + f"\n\n...[truncated, {len(content)} chars total]"
        return content

    except ImportError as e:
        return f"Missing dependency: {e}. Install: pip install pdfplumber python-docx openpyxl python-pptx"
    except Exception as e:
        return f"Document read failed: {str(e)[:200]}"


def _write_file(ws: Path, rel_path: str, content: str) -> str:
    # Protect tasks.json from direct writes
    if rel_path.strip("/") == "tasks.json":
        return "tasks.json is read-only. Use manage_tasks tool to manage tasks."

    file_path = (ws / rel_path).resolve()
    if not str(file_path).startswith(str(ws.resolve())):
        return "Access denied for this path"

    try:
        file_path.parent.mkdir(parents=True, exist_ok=True)
        file_path.write_text(content, encoding="utf-8")
        return f"✅ Written to {rel_path} ({len(content)} chars)"
    except Exception as e:
        return f"Write failed: {e}"


def _delete_file(ws: Path, rel_path: str) -> str:
    protected = {"tasks.json", "soul.md"}
    if rel_path.strip("/") in protected:
        return f"{rel_path} cannot be deleted (protected)"

    file_path = (ws / rel_path).resolve()
    if not str(file_path).startswith(str(ws.resolve())):
        return "Access denied for this path"
    if not file_path.exists():
        return f"File not found: {rel_path}"

    try:
        if file_path.is_dir():
            import shutil
            shutil.rmtree(file_path)
            return f"✅ Deleted directory {rel_path}"
        else:
            file_path.unlink()
            return f"✅ Deleted {rel_path}"
    except Exception as e:
        return f"Delete failed: {e}"


async def _manage_tasks(
    agent_id: uuid.UUID,
    user_id: uuid.UUID,
    ws: Path,
    args: dict,
) -> str:
    """Create / update / delete tasks in DB and sync to workspace."""
    from app.models.task import TaskLog
    from datetime import datetime, timezone

    action = args["action"]
    title = args["title"]

    async with async_session() as db:
        if action == "create":
            task_type = args.get("task_type", "todo")
            task = Task(
                agent_id=agent_id,
                title=title,
                description=args.get("description"),
                type=task_type,
                priority=args.get("priority", "medium"),
                created_by=user_id,
                status="pending",
                supervision_target_name=args.get("supervision_target_name"),
                supervision_channel=args.get("supervision_channel", "feishu"),
                remind_schedule=args.get("remind_schedule"),
            )
            db.add(task)
            await db.commit()
            await db.refresh(task)

            if task_type == "todo":
                # Trigger auto-execution for todo tasks
                import asyncio
                from app.services.task_executor import execute_task
                asyncio.create_task(execute_task(task.id, agent_id))
                await _sync_tasks_to_file(agent_id, ws)
                return f"✅ Task created: {title} — auto-execution started"
            else:
                # Supervision task — reminder engine will pick it up
                target = args.get('supervision_target_name', 'someone')
                schedule = args.get('remind_schedule', 'not set')
                await _sync_tasks_to_file(agent_id, ws)
                return f"✅ Supervision task created: '{title}' — will remind {target} on schedule ({schedule})"

        elif action == "update_status":
            result = await db.execute(
                select(Task).where(Task.agent_id == agent_id, Task.title.ilike(f"%{title}%"))
            )
            task = result.scalars().first()
            if not task:
                return f"No task found matching '{title}'"
            old = task.status
            task.status = args["status"]
            if args["status"] == "done":
                task.completed_at = datetime.now(timezone.utc)
            await db.commit()
            await _sync_tasks_to_file(agent_id, ws)
            return f"✅ Updated '{task.title}' from {old} to {args['status']}"

        elif action == "delete":
            from sqlalchemy import delete as sa_delete
            result = await db.execute(
                select(Task).where(Task.agent_id == agent_id, Task.title.ilike(f"%{title}%"))
            )
            task = result.scalars().first()
            if not task:
                return f"No task found matching '{title}'"
            task_title = task.title
            await db.execute(sa_delete(TaskLog).where(TaskLog.task_id == task.id))
            await db.delete(task)
            await db.commit()
            await _sync_tasks_to_file(agent_id, ws)
            return f"✅ Task deleted: {task_title}"

        return f"Unknown action: {action}"


async def _send_feishu_message(agent_id: uuid.UUID, args: dict) -> str:
    """Send a Feishu message to a person in the agent's relationship list."""
    member_name = (args.get("member_name") or "").strip()
    direct_open_id = (args.get("open_id") or "").strip()
    message_text = (args.get("message") or "").strip()

    if not message_text:
        return "❌ Please provide message content"
    if not member_name and not direct_open_id:
        return "❌ Please provide either member_name or open_id"

    try:
        from app.models.org import AgentRelationship, OrgMember
        from app.models.channel_config import ChannelConfig
        from app.services.feishu_service import feishu_service
        from sqlalchemy.orm import selectinload

        async with async_session() as db:
            # ── Shortcut: if caller already provided an open_id, use it directly ──
            if direct_open_id and not member_name:
                # Get channel config and send
                config_result = await db.execute(
                    select(ChannelConfig).where(ChannelConfig.agent_id == agent_id)
                )
                config = config_result.scalar_one_or_none()
                if not config:
                    return "❌ This agent has no Feishu channel configured"
                import json as _j
                resp = await feishu_service.send_message(
                    config.app_id, config.app_secret,
                    receive_id=direct_open_id, msg_type="text",
                    content=_j.dumps({"text": message_text}, ensure_ascii=False),
                    receive_id_type="open_id",
                )
                if resp.get("code") == 0:
                    return f"✅ 消息已发送（open_id: {direct_open_id}）"
                return f"❌ 发送失败：{resp.get('msg')} (code {resp.get('code')})"

            # Find the relationship member by name
            result = await db.execute(
                select(AgentRelationship)
                .where(AgentRelationship.agent_id == agent_id)
                .options(selectinload(AgentRelationship.member))
            )
            rels = result.scalars().all()

            target_member = None
            for r in rels:
                if r.member and r.member.name == member_name:
                    target_member = r.member
                    break

            if not target_member:
                # ── Fallback: look up via feishu_user_search (contacts cache / platform users) ──
                _search_result = await _feishu_user_search(agent_id, {"name": member_name})
                import re as _re_oid
                _oid_match = _re_oid.search(r'open_id: `(ou_[A-Za-z0-9]+)`', _search_result)
                if _oid_match:
                    _found_oid = _oid_match.group(1)
                    config_result = await db.execute(
                        select(ChannelConfig).where(ChannelConfig.agent_id == agent_id)
                    )
                    config = config_result.scalar_one_or_none()
                    if not config:
                        return "❌ This agent has no Feishu channel configured"
                    import json as _j2
                    resp = await feishu_service.send_message(
                        config.app_id, config.app_secret,
                        receive_id=_found_oid, msg_type="text",
                        content=_j2.dumps({"text": message_text}, ensure_ascii=False),
                        receive_id_type="open_id",
                    )
                    if resp.get("code") == 0:
                        return f"✅ 消息已成功发送给 {member_name}"
                    return f"❌ 找到了 {member_name}（{_found_oid}）但发送失败：{resp.get('msg')} (code {resp.get('code')})"
                # Could not find via any path
                names = [r.member.name for r in rels if r.member]
                return (
                    f"❌ 未找到联系人「{member_name}」。\n"
                    f"关系列表中的联系人：{', '.join(names) if names else '（空）'}\n"
                    f"通讯录搜索结果：{_search_result[:200]}"
                )

            if not target_member.feishu_open_id and not target_member.email and not target_member.phone:
                return f"❌ {member_name} has no linked Feishu account (no open_id, email, or phone)"

            # Get the agent's Feishu bot credentials
            config_result = await db.execute(
                select(ChannelConfig).where(ChannelConfig.agent_id == agent_id)
            )
            config = config_result.scalar_one_or_none()
            if not config:
                return "❌ This agent has no Feishu channel configured"

            import json as _json
            from app.models.system_settings import SystemSetting

            content = _json.dumps({"text": message_text}, ensure_ascii=False)

            async def _try_send(app_id: str, app_secret: str, open_id: str) -> dict:
                return await feishu_service.send_message(
                    app_id, app_secret,
                    receive_id=open_id, msg_type="text",
                    content=content, receive_id_type="open_id",
                )

            async def _save_outgoing_to_feishu_session(open_id: str):
                """Save the outgoing message to the Feishu P2P chat session."""
                try:
                    from app.models.audit import ChatMessage
                    from app.models.agent import Agent as AgentModel
                    from app.services.channel_session import find_or_create_channel_session
                    from datetime import datetime as _dt, timezone as _tz

                    agent_r = await db.execute(select(AgentModel).where(AgentModel.id == agent_id))
                    agent_obj = agent_r.scalar_one_or_none()
                    creator_id = agent_obj.creator_id if agent_obj else agent_id

                    # Look up the platform user for this Feishu open_id
                    from app.models.user import User as UserModel
                    u_r = await db.execute(
                        select(UserModel).where(UserModel.feishu_open_id == open_id)
                    )
                    feishu_user = u_r.scalar_one_or_none()
                    user_id = feishu_user.id if feishu_user else creator_id

                    ext_conv_id = f"feishu_p2p_{open_id}"
                    sess = await find_or_create_channel_session(
                        db=db,
                        agent_id=agent_id,
                        user_id=user_id,
                        external_conv_id=ext_conv_id,
                        source_channel="feishu",
                        first_message_title=f"[Agent → {member_name}]",
                    )
                    db.add(ChatMessage(
                        agent_id=agent_id,
                        user_id=user_id,
                        role="assistant",
                        content=message_text,
                        conversation_id=str(sess.id),
                    ))
                    sess.last_message_at = _dt.now(_tz.utc)
                    await db.commit()
                    print(f"[Feishu] Saved outgoing message to session {sess.id} ({member_name})")
                except Exception as e:
                    print(f"[Feishu] Failed to save outgoing message to history: {e}")

            # Step 1: Try resolve open_id for this agent's app (needs contact:user.id:readonly)
            # If successful, update stored open_id so future sends work directly
            if target_member.email or target_member.phone:
                try:
                    resolved = await feishu_service.resolve_open_id(
                        config.app_id, config.app_secret,
                        email=target_member.email,
                        mobile=target_member.phone,
                    )
                    if resolved:
                        resp = await _try_send(config.app_id, config.app_secret, resolved)
                        if resp.get("code") == 0:
                            # Update stored open_id for this app so future sends skip resolve
                            target_member.feishu_open_id = resolved
                            await db.commit()
                            await _save_outgoing_to_feishu_session(resolved)
                            return f"✅ Successfully sent message to {member_name}"
                except Exception:
                    pass

            # Step 2: Use stored open_id with agent's app (works if from same app / OAuth)
            if target_member.feishu_open_id:
                resp = await _try_send(config.app_id, config.app_secret, target_member.feishu_open_id)
                if resp.get("code") == 0:
                    await _save_outgoing_to_feishu_session(target_member.feishu_open_id)
                    return f"✅ Successfully sent message to {member_name}"

                # Step 3: If cross-app error, try org sync app as fallback
                err_msg = resp.get("msg", "")
                if "cross" in err_msg.lower():
                    org_r = await db.execute(select(SystemSetting).where(SystemSetting.key == "feishu_org_sync"))
                    org_setting = org_r.scalar_one_or_none()
                    if org_setting and org_setting.value.get("app_id"):
                        resp2 = await _try_send(
                            org_setting.value["app_id"], org_setting.value["app_secret"],
                            target_member.feishu_open_id,
                        )
                        if resp2.get("code") == 0:
                            await _save_outgoing_to_feishu_session(target_member.feishu_open_id)
                            return f"✅ Successfully sent message to {member_name}"
                        return f"❌ Send failed: {resp2.get('msg', str(resp2))}"

                return f"❌ Send failed: {err_msg}"

            return f"❌ {member_name} has no Feishu open_id and cannot be resolved via email/phone"
    except Exception as e:
        return f"❌ Message send error: {str(e)[:200]}"


async def _send_web_message(agent_id: uuid.UUID, args: dict) -> str:
    """Send a proactive message to a web platform user."""
    username = args.get("username", "").strip()
    message_text = args.get("message", "").strip()

    if not username or not message_text:
        return "❌ Please provide recipient username and message content"

    try:
        from app.models.user import User as UserModel
        from app.models.audit import ChatMessage
        from app.models.chat_session import ChatSession
        from datetime import datetime as _dt, timezone as _tz

        async with async_session() as db:
            # Look up target user by username or display_name
            from sqlalchemy import or_
            u_result = await db.execute(
                select(UserModel).where(
                    or_(
                        UserModel.username == username,
                        UserModel.display_name == username,
                    )
                )
            )
            target_user = u_result.scalar_one_or_none()
            if not target_user:
                # List available users for the agent to pick from
                all_r = await db.execute(select(UserModel.username, UserModel.display_name).limit(20))
                names = [f"{r.display_name or r.username}" for r in all_r.all()]
                return f"❌ No user named '{username}' found. Available users: {', '.join(names) if names else 'none'}"

            # Find or create a web session between the agent and this user
            sess_r = await db.execute(
                select(ChatSession).where(
                    ChatSession.agent_id == agent_id,
                    ChatSession.user_id == target_user.id,
                    ChatSession.source_channel == "web",
                ).order_by(ChatSession.created_at.desc()).limit(1)
            )
            session = sess_r.scalar_one_or_none()

            if not session:
                # Create a new session for this user
                session = ChatSession(
                    agent_id=agent_id,
                    user_id=target_user.id,
                    title=f"[Agent Message] {_dt.now(_tz.utc).strftime('%m-%d %H:%M')}",
                    source_channel="web",
                    created_at=_dt.now(_tz.utc),
                )
                db.add(session)
                await db.flush()

            # Save the message
            db.add(ChatMessage(
                agent_id=agent_id,
                user_id=target_user.id,
                role="assistant",
                content=message_text,
                conversation_id=str(session.id),
            ))
            session.last_message_at = _dt.now(_tz.utc)
            await db.commit()

            # Push via WebSocket if user has an active connection
            try:
                from app.api.websocket import manager as ws_manager
                agent_id_str = str(agent_id)
                if agent_id_str in ws_manager.active_connections:
                    for ws, sid in list(ws_manager.active_connections[agent_id_str]):
                        try:
                            await ws.send_json({
                                "type": "trigger_notification",
                                "content": message_text,
                                "triggers": ["web_message"],
                            })
                        except Exception:
                            pass
            except Exception:
                pass

            display = target_user.display_name or target_user.username
            return f"✅ Message sent to {display} on web platform. It has been saved to their chat history."

    except Exception as e:
        return f"❌ Web message send error: {str(e)[:200]}"


async def _send_message_to_agent(from_agent_id: uuid.UUID, args: dict) -> str:
    """Send a message to another digital employee. Uses a single request-response pattern:
    the source agent sends a message, the target agent replies once, and the result is returned.
    If the source agent needs to continue the conversation, it can call this tool again.
    """
    agent_name = args.get("agent_name", "").strip()
    message_text = args.get("message", "").strip()

    if not agent_name or not message_text:
        return "❌ Please provide target agent name and message content"

    try:
        from app.models.agent import Agent
        from app.models.audit import ChatMessage
        from app.models.chat_session import ChatSession
        from app.models.participant import Participant
        from datetime import datetime, timezone

        async with async_session() as db:
            # Look up source agent
            src_result = await db.execute(select(Agent).where(Agent.id == from_agent_id))
            source_agent = src_result.scalar_one_or_none()
            source_name = source_agent.name if source_agent else "Unknown agent"

            # Find target agent by name
            result = await db.execute(
                select(Agent).where(Agent.name.ilike(f"%{agent_name}%"), Agent.id != from_agent_id)
            )
            target = result.scalars().first()
            if not target:
                all_r = await db.execute(select(Agent).where(Agent.id != from_agent_id))
                names = [a.name for a in all_r.scalars().all()]
                return f"❌ No agent found matching '{agent_name}'. Available: {', '.join(names) if names else 'none'}"

            # Check if target agent has expired
            if target.is_expired or (target.expires_at and datetime.now(timezone.utc) >= target.expires_at):
                return f"⚠️ {target.name} is currently unavailable — their service period has ended. Please contact the platform administrator."

            # Get Participant IDs for both agents
            src_part_r = await db.execute(select(Participant).where(Participant.type == "agent", Participant.ref_id == from_agent_id))
            src_participant = src_part_r.scalar_one_or_none()
            tgt_part_r = await db.execute(select(Participant).where(Participant.type == "agent", Participant.ref_id == target.id))
            tgt_participant = tgt_part_r.scalar_one_or_none()

            # Find or create ChatSession for this agent pair (ordered consistently)
            session_agent_id = min(from_agent_id, target.id, key=str)
            session_peer_id = max(from_agent_id, target.id, key=str)
            sess_r = await db.execute(
                select(ChatSession).where(
                    ChatSession.agent_id == session_agent_id,
                    ChatSession.peer_agent_id == session_peer_id,
                    ChatSession.source_channel == "agent",
                )
            )
            chat_session = sess_r.scalar_one_or_none()
            if not chat_session:
                owner_id = source_agent.creator_id if source_agent else from_agent_id
                src_part_id = src_participant.id if src_participant else None
                chat_session = ChatSession(
                    agent_id=session_agent_id,
                    user_id=owner_id,
                    title=f"{source_name} ↔ {target.name}",
                    source_channel="agent",
                    participant_id=src_part_id,
                    peer_agent_id=session_peer_id,
                )
                db.add(chat_session)
                await db.flush()

            session_id = str(chat_session.id)

            # Prepare target LLM
            from app.services.agent_context import build_agent_context
            from app.models.llm import LLMModel
            import httpx

            if not target.primary_model_id:
                return f"⚠️ {target.name} has no LLM model configured"

            model_r = await db.execute(select(LLMModel).where(LLMModel.id == target.primary_model_id))
            target_model = model_r.scalar_one_or_none()
            if not target_model:
                return f"⚠️ {target.name}'s model configuration is invalid"

            # Build target system prompt
            target_system = await build_agent_context(target.id, target.name, target.role_description or "")
            target_system += (
                "\n\n--- Agent-to-Agent Message ---\n"
                "You are receiving a message from another digital employee. "
                "Reply concisely and helpfully. Focus on the request and provide a clear answer.\n"
            )

            # Load recent history for context
            conversation_messages: list[dict] = []
            hist_result = await db.execute(
                select(ChatMessage)
                .where(
                    ChatMessage.conversation_id == session_id,
                    ChatMessage.agent_id == session_agent_id,
                )
                .order_by(ChatMessage.created_at.desc())
                .limit(6)
            )
            for m in reversed(hist_result.scalars().all()):
                if m.participant_id and src_participant and m.participant_id == src_participant.id:
                    role = "user"
                else:
                    role = "assistant"
                conversation_messages.append({"role": role, "content": m.content})

            # Add the new message from source
            conversation_messages.append({"role": "user", "content": f"[From {source_name}] {message_text}"})

            # Save source message
            owner_id = source_agent.creator_id if source_agent else from_agent_id
            db.add(ChatMessage(
                agent_id=session_agent_id,
                user_id=owner_id,
                role="user",
                content=message_text,
                conversation_id=session_id,
                participant_id=src_participant.id if src_participant else None,
            ))
            chat_session.last_message_at = datetime.now(timezone.utc)
            await db.commit()

            # Call target LLM with tool support (multi-round)
            from app.services.llm_utils import get_provider_base_url
            from app.services.agent_tools import get_agent_tools_for_llm, execute_tool
            base_url = get_provider_base_url(target_model.provider, target_model.base_url)
            url = f"{base_url.rstrip('/')}/chat/completions"
            full_msgs = [{"role": "system", "content": target_system}] + conversation_messages

            # Load tools for target agent
            tools_for_llm = await get_agent_tools_for_llm(target.id)

            max_tool_rounds = 5
            target_reply = ""
            _a2a_accumulated_tokens = 0

            from app.services.token_tracker import record_token_usage, extract_usage_tokens, estimate_tokens_from_chars

            async with httpx.AsyncClient(timeout=120) as client:
                for _round in range(max_tool_rounds):
                    req_body: dict = {
                        "model": target_model.model,
                        "messages": full_msgs,
                        "temperature": 0.7,
                        "max_tokens": 4096,
                    }
                    if tools_for_llm:
                        req_body["tools"] = tools_for_llm
                        req_body["tool_choice"] = "auto"

                    resp = await client.post(
                        url,
                        json=req_body,
                        headers={"Authorization": f"Bearer {target_model.api_key_encrypted}"},
                    )
                    data = resp.json()

                    # Track tokens from API response
                    real_tokens = extract_usage_tokens(data.get("usage"))
                    if real_tokens:
                        _a2a_accumulated_tokens += real_tokens
                    else:
                        round_chars = sum(len(m.get("content") or '') for m in full_msgs if isinstance(m.get("content"), str))
                        _a2a_accumulated_tokens += estimate_tokens_from_chars(round_chars)

                    if "choices" not in data or not data["choices"]:
                        break

                    choice = data["choices"][0]
                    msg = choice.get("message", {})

                    # Check for tool calls
                    tool_calls = msg.get("tool_calls")
                    if tool_calls:
                        # Add assistant message with tool calls to conversation
                        full_msgs.append(msg)

                        # Execute each tool call
                        for tc in tool_calls:
                            fn = tc.get("function", {})
                            tool_name = fn.get("name", "")
                            try:
                                tool_args = json.loads(fn.get("arguments", "{}"))
                            except Exception:
                                tool_args = {}

                            tool_result = await execute_tool(tool_name, tool_args, target.id, owner_id)

                            # Add tool result to conversation
                            full_msgs.append({
                                "role": "tool",
                                "tool_call_id": tc.get("id", ""),
                                "content": tool_result[:4000],
                            })
                        continue  # Next LLM round

                    # No tool calls — this is the final text response
                    target_reply = msg.get("content", "")
                    break

            # Record accumulated A2A tokens for the target agent
            if _a2a_accumulated_tokens > 0:
                await record_token_usage(target.id, _a2a_accumulated_tokens)

            if not target_reply:
                return f"⚠️ {target.name} did not respond (LLM returned empty)"

            # Save target reply
            async with async_session() as db2:
                part_r = await db2.execute(select(Participant).where(Participant.type == "agent", Participant.ref_id == target.id))
                tgt_part = part_r.scalar_one_or_none()
                db2.add(ChatMessage(
                    agent_id=session_agent_id,
                    user_id=owner_id,
                    role="assistant",
                    content=target_reply,
                    conversation_id=session_id,
                    participant_id=tgt_part.id if tgt_part else None,
                ))
                await db2.commit()

            # Log activity
            from app.services.activity_logger import log_activity
            await log_activity(
                target.id, "agent_msg_sent",
                f"Replied to message from {source_name}",
                detail={"partner": source_name, "message": message_text[:200], "reply": target_reply[:200]},
            )
            await log_activity(
                from_agent_id, "agent_msg_sent",
                f"Sent message to {target.name} and received reply",
                detail={"partner": target.name, "message": message_text[:200], "reply": target_reply[:200]},
            )

            return f"💬 {target.name} replied:\n{target_reply}"

    except Exception as e:
        import traceback
        traceback.print_exc()
        return f"❌ Message send error: {str(e)[:200]}"



# Plaza Tools — Agent Square social feed
# ═══════════════════════════════════════════════════════

async def _plaza_get_new_posts(arguments: dict) -> str:
    """Get recent posts from the Agent Plaza."""
    from app.models.plaza import PlazaPost, PlazaComment
    from sqlalchemy import desc

    limit = min(arguments.get("limit", 10), 20)

    try:
        async with async_session() as db:
            q = select(PlazaPost).order_by(desc(PlazaPost.created_at)).limit(limit)
            result = await db.execute(q)
            posts = result.scalars().all()

            if not posts:
                return "📭 No posts in the plaza yet. Be the first to share something!"

            output = []
            for p in posts:
                # Load comments
                cr = await db.execute(
                    select(PlazaComment).where(PlazaComment.post_id == p.id).order_by(PlazaComment.created_at).limit(5)
                )
                comments = cr.scalars().all()
                icon = "🤖" if p.author_type == "agent" else "👤"
                time_str = p.created_at.strftime("%m-%d %H:%M") if p.created_at else ""
                post_text = f"{icon} **{p.author_name}** ({time_str}) [post_id: {p.id}]\n{p.content}\n❤️ {p.likes_count}  💬 {p.comments_count}"
                if comments:
                    for c in comments:
                        c_icon = "🤖" if c.author_type == "agent" else "👤"
                        post_text += f"\n  └─ {c_icon} {c.author_name}: {c.content}"
                output.append(post_text)

            return "🏛️ Agent Plaza — Recent Posts:\n\n" + "\n\n---\n\n".join(output)

    except Exception as e:
        return f"❌ Failed to load plaza posts: {str(e)[:200]}"


async def _plaza_create_post(agent_id: uuid.UUID, arguments: dict) -> str:
    """Create a new post in the Agent Plaza."""
    from app.models.plaza import PlazaPost
    from app.models.agent import Agent as AgentModel

    content = arguments.get("content", "").strip()
    if not content:
        return "❌ Post content cannot be empty."
    if len(content) > 500:
        content = content[:500]

    try:
        async with async_session() as db:
            # Get agent name
            ar = await db.execute(select(AgentModel).where(AgentModel.id == agent_id))
            agent = ar.scalar_one_or_none()
            if not agent:
                return "❌ Agent not found."

            post = PlazaPost(
                author_id=agent_id,
                author_type="agent",
                author_name=agent.name,
                content=content,
                tenant_id=agent.tenant_id,
            )
            db.add(post)
            await db.commit()
            await db.refresh(post)
            return f"✅ Post published! (ID: {post.id})"

    except Exception as e:
        return f"❌ Failed to create post: {str(e)[:200]}"


async def _plaza_add_comment(agent_id: uuid.UUID, arguments: dict) -> str:
    """Add a comment to a plaza post."""
    from app.models.plaza import PlazaPost, PlazaComment
    from app.models.agent import Agent as AgentModel

    post_id = arguments.get("post_id", "")
    content = arguments.get("content", "").strip()
    if not content:
        return "❌ Comment content cannot be empty."
    if len(content) > 300:
        content = content[:300]

    try:
        pid = uuid.UUID(str(post_id))
    except Exception:
        return "❌ Invalid post_id format."

    try:
        async with async_session() as db:
            # Verify post exists
            pr = await db.execute(select(PlazaPost).where(PlazaPost.id == pid))
            post = pr.scalar_one_or_none()
            if not post:
                return "❌ Post not found."

            # Get agent name
            ar = await db.execute(select(AgentModel).where(AgentModel.id == agent_id))
            agent = ar.scalar_one_or_none()
            if not agent:
                return "❌ Agent not found."

            comment = PlazaComment(
                post_id=pid,
                author_id=agent_id,
                author_type="agent",
                author_name=agent.name,
                content=content,
            )
            db.add(comment)
            post.comments_count = (post.comments_count or 0) + 1
            await db.commit()
            return f"✅ Comment added to post by {post.author_name}."

    except Exception as e:
        return f"❌ Failed to add comment: {str(e)[:200]}"


# ─── Code Execution ─────────────────────────────────────────────

# Dangerous patterns to block
_DANGEROUS_BASH = [
    "rm -rf /", "rm -rf ~", "sudo ", "mkfs", "dd if=",
    ":(){ :", "chmod 777 /", "chown ", "shutdown", "reboot",
    "curl ", "wget ", "nc ", "ncat ", "ssh ", "scp ",
    "python3 -c", "python -c",
]

_DANGEROUS_PYTHON_IMPORTS = [
    "subprocess", "shutil.rmtree", "os.system", "os.popen",
    "os.exec", "os.spawn",
    "socket", "http.client", "urllib.request", "requests",
    "ftplib", "smtplib", "telnetlib", "ctypes",
    "__import__", "importlib",
]


def _check_code_safety(language: str, code: str) -> str | None:
    """Check code for dangerous patterns. Returns error message if unsafe, None if ok."""
    code_lower = code.lower()

    if language == "bash":
        for pattern in _DANGEROUS_BASH:
            if pattern.lower() in code_lower:
                return f"❌ Blocked: dangerous command detected ({pattern.strip()})"
        # Block deep path traversal outside workspace
        if "../../" in code:
            return "❌ Blocked: directory traversal not allowed"

    elif language == "python":
        for pattern in _DANGEROUS_PYTHON_IMPORTS:
            if pattern.lower() in code_lower:
                return f"❌ Blocked: unsafe operation detected ({pattern})"

    elif language == "node":
        dangerous_node = ["child_process", "fs.rmSync", "fs.rmdirSync", "process.exit",
                          "require('http')", "require('https')", "require('net')"]
        for pattern in dangerous_node:
            if pattern.lower() in code_lower:
                return f"❌ Blocked: unsafe operation detected ({pattern})"

    return None


async def _execute_code(ws: Path, arguments: dict) -> str:
    """Execute code in a sandboxed subprocess within the agent's workspace."""
    import asyncio

    language = arguments.get("language", "python")
    code = arguments.get("code", "")
    timeout = min(arguments.get("timeout", 30), 60)  # Max 60 seconds

    if not code.strip():
        return "❌ No code provided"

    if language not in ("python", "bash", "node"):
        return f"❌ Unsupported language: {language}. Use: python, bash, or node"

    # Security check
    safety_error = _check_code_safety(language, code)
    if safety_error:
        return safety_error

    # Working directory is the agent's workspace/ subdirectory (must be absolute)
    work_dir = (ws / "workspace").resolve()
    work_dir.mkdir(parents=True, exist_ok=True)

    # Determine command and file extension
    if language == "python":
        ext = ".py"
        cmd_prefix = ["python3"]
    elif language == "bash":
        ext = ".sh"
        cmd_prefix = ["bash"]
    elif language == "node":
        ext = ".js"
        cmd_prefix = ["node"]
    else:
        return f"❌ Unsupported language: {language}"

    # Write code to a temp file inside workspace
    script_path = work_dir / f"_exec_tmp{ext}"
    try:
        script_path.write_text(code, encoding="utf-8")

        # Inherit parent environment but override HOME to workspace
        safe_env = dict(os.environ)
        safe_env["HOME"] = str(work_dir)
        safe_env["PYTHONDONTWRITEBYTECODE"] = "1"

        proc = await asyncio.create_subprocess_exec(
            *cmd_prefix, str(script_path),
            cwd=str(work_dir),
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.PIPE,
            env=safe_env,
        )

        try:
            stdout, stderr = await asyncio.wait_for(proc.communicate(), timeout=timeout)
        except asyncio.TimeoutError:
            proc.kill()
            await proc.communicate()
            return f"❌ Code execution timed out after {timeout}s"

        stdout_str = stdout.decode("utf-8", errors="replace")[:10000]
        stderr_str = stderr.decode("utf-8", errors="replace")[:5000]

        result_parts = []
        if stdout_str.strip():
            result_parts.append(f"📤 Output:\n{stdout_str}")
        if stderr_str.strip():
            result_parts.append(f"⚠️ Stderr:\n{stderr_str}")
        if proc.returncode != 0:
            result_parts.append(f"Exit code: {proc.returncode}")

        if not result_parts:
            return "✅ Code executed successfully (no output)"

        return "\n\n".join(result_parts)

    except Exception as e:
        return f"❌ Execution error: {str(e)[:200]}"
    finally:
        # Clean up temp script
        try:
            script_path.unlink(missing_ok=True)
        except Exception:
            pass


# ─── Resource Discovery Executors ───────────────────────────────

async def _discover_resources(arguments: dict) -> str:
    """Search Smithery registry for MCP servers."""
    query = arguments.get("query", "")
    if not query:
        return "❌ Please provide a search query describing the capability you need."
    max_results = min(arguments.get("max_results", 5), 10)

    from app.services.resource_discovery import search_smithery
    return await search_smithery(query, max_results)


async def _import_mcp_server(agent_id: uuid.UUID, arguments: dict) -> str:
    """Import an MCP server — either from Smithery or by direct URL."""
    config = arguments.get("config") or {}
    reauthorize = arguments.get("reauthorize", False)
    mcp_url = config.pop("mcp_url", None) if isinstance(config, dict) else None

    if mcp_url:
        # Direct URL import — bypass Smithery
        from app.services.resource_discovery import import_mcp_direct
        server_name = arguments.get("server_id") or config.pop("server_name", None)
        api_key = config.pop("api_key", None)
        return await import_mcp_direct(mcp_url, agent_id, server_name, api_key)

    # Smithery import
    server_id = arguments.get("server_id", "")
    if not server_id:
        return "❌ Please provide a server_id (e.g. 'github'). Use discover_resources first to find available servers."

    from app.services.resource_discovery import import_mcp_from_smithery
    return await import_mcp_from_smithery(server_id, agent_id, config or None, reauthorize=reauthorize)


# ─── Trigger Management Handlers (Pulse Engine) ────────────────────

MAX_TRIGGERS_PER_AGENT = 20
VALID_TRIGGER_TYPES = {"cron", "once", "interval", "poll", "on_message", "webhook"}


async def _handle_set_trigger(agent_id: uuid.UUID, arguments: dict) -> str:
    """Create a new trigger for the agent."""
    from app.models.trigger import AgentTrigger

    name = arguments.get("name", "").strip()
    ttype = arguments.get("type", "").strip()
    config = arguments.get("config", {})
    reason = arguments.get("reason", "").strip()
    focus_ref = arguments.get("focus_ref", "") or arguments.get("agenda_ref", "")  # backward compat

    if not name:
        return "❌ Missing required argument 'name'"
    if ttype not in VALID_TRIGGER_TYPES:
        return f"❌ Invalid trigger type '{ttype}'. Valid types: {', '.join(VALID_TRIGGER_TYPES)}"
    if not reason:
        return "❌ Missing required argument 'reason'"

    # Validate type-specific config
    if ttype == "cron":
        expr = config.get("expr", "")
        if not expr:
            return "❌ cron trigger requires config.expr, e.g. {\"expr\": \"0 9 * * *\"}"
        try:
            from croniter import croniter
            croniter(expr)
        except Exception:
            return f"❌ Invalid cron expression: '{expr}'"
    elif ttype == "once":
        if not config.get("at"):
            return "❌ once trigger requires config.at, e.g. {\"at\": \"2026-03-10T09:00:00+08:00\"}"
    elif ttype == "interval":
        if not config.get("minutes"):
            return "❌ interval trigger requires config.minutes, e.g. {\"minutes\": 30}"
    elif ttype == "poll":
        if not config.get("url"):
            return "❌ poll trigger requires config.url"
    elif ttype == "on_message":
        if not config.get("from_agent_name") and not config.get("from_user_name"):
            return "❌ on_message trigger requires config.from_agent_name (for agents) or config.from_user_name (for human users on Feishu/Slack/Discord)"
    elif ttype == "webhook":
        # Auto-generate a unique token for the webhook URL
        import secrets
        token = secrets.token_urlsafe(8)  # ~11 chars, URL-safe
        config["token"] = token

    try:
        async with async_session() as db:
            # Load agent to get per-agent trigger limit
            from app.models.agent import Agent as _AgentModel
            _a_result = await db.execute(select(_AgentModel).where(_AgentModel.id == agent_id))
            _agent_obj = _a_result.scalar_one_or_none()
            agent_max_triggers = (_agent_obj.max_triggers if _agent_obj else None) or MAX_TRIGGERS_PER_AGENT

            # Check max triggers
            from sqlalchemy import func as sa_func
            result = await db.execute(
                select(sa_func.count()).select_from(AgentTrigger).where(
                    AgentTrigger.agent_id == agent_id,
                    AgentTrigger.is_enabled == True,
                )
            )
            count = result.scalar() or 0
            if count >= agent_max_triggers:
                return f"❌ Maximum trigger limit reached ({agent_max_triggers}). Cancel some triggers first."

            # Check for duplicate name
            result = await db.execute(
                select(AgentTrigger).where(
                    AgentTrigger.agent_id == agent_id,
                    AgentTrigger.name == name,
                )
            )
            existing = result.scalar_one_or_none()
            if existing:
                if existing.is_enabled:
                    return f"❌ Trigger '{name}' already exists and is active. Use update_trigger to modify it, or cancel_trigger first."
                else:
                    # Re-enable disabled trigger with new config
                    existing.type = ttype
                    existing.config = config
                    existing.reason = reason
                    existing.focus_ref = focus_ref or None
                    existing.is_enabled = True
                    existing.fire_count = 0
                    existing.last_fired_at = None
                    await db.commit()
                    return f"✅ Trigger '{name}' re-enabled with new configuration ({ttype})"

            trigger = AgentTrigger(
                agent_id=agent_id,
                name=name,
                type=ttype,
                config=config,
                reason=reason,
                focus_ref=focus_ref or None,
            )
            db.add(trigger)
            await db.commit()

        # Activity log
        try:
            from app.services.audit_logger import write_audit_log
            await write_audit_log("trigger_created", {
                "name": name, "type": ttype, "reason": reason[:100],
            }, agent_id=agent_id)
        except Exception:
            pass

        # Return webhook URL for webhook triggers
        if ttype == "webhook":
            from app.config import get_settings
            settings = get_settings()
            base = getattr(settings, 'PUBLIC_URL', '') or ''
            if not base:
                base = 'https://try.clawith.ai'  # fallback
            webhook_url = f"{base.rstrip('/')}/api/webhooks/t/{config['token']}"
            return f"✅ Webhook trigger '{name}' created.\n\nWebhook URL: {webhook_url}\n\nTell the user to configure this URL in their external service (e.g. GitHub, Grafana). When the service sends a POST to this URL, you will be woken up with the payload as context."

        return f"✅ Trigger '{name}' created ({ttype}). It will fire according to your config and wake you up with the reason as context."

    except Exception as e:
        return f"❌ Failed to create trigger: {e}"


async def _handle_update_trigger(agent_id: uuid.UUID, arguments: dict) -> str:
    """Update an existing trigger's config or reason."""
    from app.models.trigger import AgentTrigger

    name = arguments.get("name", "").strip()
    if not name:
        return "❌ Missing required argument 'name'"

    new_config = arguments.get("config")
    new_reason = arguments.get("reason")

    if new_config is None and new_reason is None:
        return "❌ Provide at least one of 'config' or 'reason' to update"

    try:
        async with async_session() as db:
            result = await db.execute(
                select(AgentTrigger).where(
                    AgentTrigger.agent_id == agent_id,
                    AgentTrigger.name == name,
                )
            )
            trigger = result.scalar_one_or_none()
            if not trigger:
                return f"❌ Trigger '{name}' not found"

            changes = []
            if new_config is not None:
                old_config = trigger.config
                trigger.config = new_config
                changes.append(f"config: {old_config} → {new_config}")
            if new_reason is not None:
                trigger.reason = new_reason
                changes.append(f"reason updated")

            await db.commit()

        try:
            from app.services.audit_logger import write_audit_log
            await write_audit_log("trigger_updated", {
                "name": name, "changes": "; ".join(changes),
            }, agent_id=agent_id)
        except Exception:
            pass

        return f"✅ Trigger '{name}' updated: {'; '.join(changes)}"

    except Exception as e:
        return f"❌ Failed to update trigger: {e}"


async def _handle_cancel_trigger(agent_id: uuid.UUID, arguments: dict) -> str:
    """Cancel (disable) a trigger by name."""
    from app.models.trigger import AgentTrigger

    name = arguments.get("name", "").strip()
    if not name:
        return "❌ Missing required argument 'name'"

    try:
        async with async_session() as db:
            result = await db.execute(
                select(AgentTrigger).where(
                    AgentTrigger.agent_id == agent_id,
                    AgentTrigger.name == name,
                )
            )
            trigger = result.scalar_one_or_none()
            if not trigger:
                return f"❌ Trigger '{name}' not found"
            if not trigger.is_enabled:
                return f"ℹ️ Trigger '{name}' is already disabled"

            trigger.is_enabled = False
            await db.commit()

        try:
            from app.services.audit_logger import write_audit_log
            await write_audit_log("trigger_cancelled", {"name": name}, agent_id=agent_id)
        except Exception:
            pass

        return f"✅ Trigger '{name}' cancelled. It will no longer fire."

    except Exception as e:
        return f"❌ Failed to cancel trigger: {e}"


async def _handle_list_triggers(agent_id: uuid.UUID) -> str:
    """List all active triggers for the agent."""
    from app.models.trigger import AgentTrigger

    try:
        async with async_session() as db:
            result = await db.execute(
                select(AgentTrigger).where(
                    AgentTrigger.agent_id == agent_id,
                ).order_by(AgentTrigger.created_at.desc())
            )
            triggers = result.scalars().all()

        if not triggers:
            return "No triggers found. Use set_trigger to create one."

        lines = ["| Name | Type | Config | Reason | Status | Fires |", "|------|------|--------|--------|--------|-------|"]
        for t in triggers:
            status = "✅ active" if t.is_enabled else "⏸ disabled"
            config_str = str(t.config)[:50]
            reason_str = t.reason[:40] if t.reason else ""
            lines.append(f"| {t.name} | {t.type} | {config_str} | {reason_str} | {status} | {t.fire_count} |")

        return "\n".join(lines)

    except Exception as e:
        return f"❌ Failed to list triggers: {e}"


# ─── Image Upload (ImageKit CDN) ────────────────────────────────

async def _upload_image(agent_id: uuid.UUID, ws: Path, arguments: dict) -> str:
    """Upload an image to ImageKit CDN and return the public URL.

    Credential resolution order:
    1. Global tool config (admin-set, shared by all agents)
    2. Per-agent tool config override (agent-specific)
    """
    import httpx
    import base64

    file_path = arguments.get("file_path")
    url = arguments.get("url")
    file_name = arguments.get("file_name")
    folder = arguments.get("folder", "/clawith")

    if not file_path and not url:
        return "❌ Please provide either 'file_path' (workspace path) or 'url' (public image URL)"

    # ── Load ImageKit credentials (global → per-agent fallback) ──
    private_key = ""
    url_endpoint = ""
    try:
        from app.models.tool import Tool, AgentTool
        async with async_session() as db:
            # Global config
            r = await db.execute(select(Tool).where(Tool.name == "upload_image"))
            tool = r.scalar_one_or_none()
            if tool and tool.config:
                private_key = tool.config.get("private_key", "")
                url_endpoint = tool.config.get("url_endpoint", "")

            # Per-agent override (if global key is empty)
            if not private_key and tool:
                r2 = await db.execute(
                    select(AgentTool).where(
                        AgentTool.agent_id == agent_id,
                        AgentTool.tool_id == tool.id,
                    )
                )
                agent_tool = r2.scalar_one_or_none()
                if agent_tool and agent_tool.config:
                    private_key = agent_tool.config.get("private_key", "") or private_key
                    url_endpoint = agent_tool.config.get("url_endpoint", "") or url_endpoint
    except Exception as e:
        print(f"[UploadImage] Config load error: {e}")

    if not private_key:
        return "❌ ImageKit Private Key not configured. Ask your admin to configure it in Enterprise Settings → Tools → Upload Image, or set it in your agent's tool config."

    # ── Prepare the file ──
    form_data = {}
    file_content = None

    if file_path:
        # Read from workspace
        full_path = (ws / file_path).resolve()
        if not str(full_path).startswith(str(ws)):
            return "❌ Access denied: path is outside the workspace"
        if not full_path.exists():
            return f"❌ File not found: {file_path}"
        if not full_path.is_file():
            return f"❌ Not a file: {file_path}"

        # Check file size (max 25MB for free plan)
        size_mb = full_path.stat().st_size / (1024 * 1024)
        if size_mb > 25:
            return f"❌ File too large ({size_mb:.1f}MB). Maximum is 25MB."

        file_content = full_path.read_bytes()
        if not file_name:
            file_name = full_path.name
    elif url:
        # Pass URL directly to ImageKit
        form_data["file"] = url
        if not file_name:
            from urllib.parse import urlparse
            file_name = urlparse(url).path.split("/")[-1] or "image.jpg"

    if not file_name:
        file_name = "image.png"

    form_data["fileName"] = file_name
    form_data["folder"] = folder
    form_data["useUniqueFileName"] = "true"

    # ── Upload to ImageKit V2 ──
    auth_string = base64.b64encode(f"{private_key}:".encode()).decode()

    try:
        async with httpx.AsyncClient(timeout=60) as client:
            if file_content:
                # Binary upload via multipart
                files = {"file": (file_name, file_content)}
                resp = await client.post(
                    "https://upload.imagekit.io/api/v2/files/upload",
                    headers={"Authorization": f"Basic {auth_string}"},
                    data=form_data,
                    files=files,
                )
            else:
                # URL upload via form data
                resp = await client.post(
                    "https://upload.imagekit.io/api/v2/files/upload",
                    headers={"Authorization": f"Basic {auth_string}"},
                    data=form_data,
                )

        if resp.status_code in (200, 201):
            result = resp.json()
            cdn_url = result.get("url", "")
            file_id = result.get("fileId", "")
            size = result.get("size", 0)
            size_str = f"{size / 1024:.1f}KB" if size < 1024 * 1024 else f"{size / (1024 * 1024):.1f}MB"
            return (
                f"✅ Image uploaded successfully!\n\n"
                f"**CDN URL**: {cdn_url}\n"
                f"**File ID**: {file_id}\n"
                f"**Size**: {size_str}\n"
                f"**Name**: {result.get('name', file_name)}"
            )
        else:
            error_detail = resp.text[:300]
            return f"❌ Upload failed (HTTP {resp.status_code}): {error_detail}"

    except httpx.TimeoutException:
        return "❌ Upload timed out after 60s. The file may be too large or the network is slow."
    except Exception as e:
        return f"❌ Upload error: {type(e).__name__}: {str(e)[:300]}"



# ─── Feishu Helper ────────────────────────────────────────────────────────────

async def _get_feishu_token(agent_id: uuid.UUID) -> tuple[str, str] | None:
    """Get (app_id, app_access_token) for the agent's configured Feishu channel."""
    import httpx
    from app.models.channel_config import ChannelConfig

    async with async_session() as db:
        result = await db.execute(
            select(ChannelConfig).where(
                ChannelConfig.agent_id == agent_id,
                ChannelConfig.channel_type == "feishu",
                ChannelConfig.is_configured == True,
            )
        )
        config = result.scalar_one_or_none()

    if not config or not config.app_id or not config.app_secret:
        return None

    async with httpx.AsyncClient(timeout=10) as client:
        resp = await client.post(
            "https://open.feishu.cn/open-apis/auth/v3/tenant_access_token/internal",
            json={"app_id": config.app_id, "app_secret": config.app_secret},
        )
        token = resp.json().get("tenant_access_token", "")

    return (config.app_id, token) if token else None


async def _get_agent_calendar_id(token: str) -> tuple[str | None, str | None]:
    """Get (calendar_id, error_msg) for the agent app's primary calendar.

    Returns (calendar_id, None) on success, or (None, human_readable_error) on failure.
    """
    import httpx
    async with httpx.AsyncClient(timeout=10) as client:
        resp = await client.post(
            "https://open.feishu.cn/open-apis/calendar/v4/calendars/primary",
            headers={"Authorization": f"Bearer {token}"},
        )
    data = resp.json()
    code = data.get("code", -1)
    if code == 0:
        cals = data.get("data", {}).get("calendars", [])
        if cals:
            cal_id = cals[0].get("calendar", {}).get("calendar_id")
            return cal_id, None
        return None, "日历列表为空，请确认应用有 calendar:calendar 权限并已发布新版本"
    if code == 99991672:
        return None, (
            "❌ 飞书日历权限未开通（错误码 99991672）\n\n"
            "请在飞书开放平台为应用 cli_a9257c5136781ceb 开通以下权限并发布新版本：\n"
            "• calendar:calendar:readonly（应用身份权限）\n"
            "• calendar:calendar.event:create（应用身份权限）\n"
            "• calendar:calendar.event:read（用户身份权限）\n"
            "• calendar:calendar.event:update（用户身份权限）\n"
            "• calendar:calendar.event:delete（用户身份权限）\n\n"
            "开通步骤：飞书开放平台 → 权限管理 → 批量导入权限 → 添加以上权限 → 创建版本 → 确认发布"
        )
    return None, f"获取日历 ID 失败：{data.get('msg')} (code {code})"


async def _feishu_resolve_open_id(token: str, email: str) -> str | None:
    """Resolve a user's open_id from their email."""
    import httpx
    async with httpx.AsyncClient(timeout=10) as client:
        resp = await client.post(
            "https://open.feishu.cn/open-apis/contact/v3/users/batch_get_id",
            json={"emails": [email]},
            headers={"Authorization": f"Bearer {token}"},
            params={"user_id_type": "open_id"},
        )
    data = resp.json()
    if data.get("code") != 0:
        return None
    for u in data.get("data", {}).get("user_list", []):
        oid = u.get("user_id")
        if oid:
            return oid
    return None


def _iso_to_ts(iso_str: str) -> float:
    """Convert ISO 8601 string to Unix timestamp."""
    from datetime import datetime as _dt
    for fmt in ("%Y-%m-%dT%H:%M:%S%z", "%Y-%m-%dT%H:%M:%SZ", "%Y-%m-%dT%H:%M:%S"):
        try:
            if iso_str.endswith("Z"):
                d = _dt.fromisoformat(iso_str.replace("Z", "+00:00"))
            else:
                d = _dt.strptime(iso_str, fmt)
            return d.timestamp()
        except ValueError:
            continue
    raise ValueError(f"Cannot parse datetime: {iso_str!r}")


# ─── Feishu Document Tools ────────────────────────────────────────────────────

# ─── Feishu Wiki Tools ───────────────────────────────────────────────────────

async def _feishu_wiki_get_node(token_str: str, auth_token: str) -> dict | None:
    """Call wiki get_node API to resolve a wiki node token → {obj_token, space_id, has_child, title}.
    Returns None if the token is not a wiki node."""
    import httpx
    async with httpx.AsyncClient(timeout=5) as client:
        r = await client.get(
            "https://open.feishu.cn/open-apis/wiki/v2/spaces/get_node",
            headers={"Authorization": f"Bearer {auth_token}"},
            params={"token": token_str, "obj_type": "wiki"},
        )
    d = r.json()
    if d.get("code") != 0:
        return None
    node = d.get("data", {}).get("node", {})
    return {
        "obj_token": node.get("obj_token", ""),
        "space_id": node.get("origin_space_id", node.get("space_id", "")),
        "has_child": node.get("has_child", False),
        "title": node.get("title", ""),
        "node_token": node.get("node_token", token_str),
    }


async def _feishu_wiki_list(agent_id: uuid.UUID, arguments: dict) -> str:
    """List sub-pages of a Feishu Wiki node, optionally recursive."""
    import httpx

    node_token = (arguments.get("node_token") or "").strip()
    recursive = bool(arguments.get("recursive", False))

    if not node_token:
        return "❌ Missing required argument 'node_token'"

    creds = await _get_feishu_token(agent_id)
    if not creds:
        return "❌ Agent has no Feishu channel configured."
    _, token = creds
    headers = {"Authorization": f"Bearer {token}"}

    # Resolve node → space_id
    node_info = await _feishu_wiki_get_node(node_token, token)
    if not node_info:
        return (
            f"❌ 无法解析 Wiki 节点 `{node_token}`。\n"
            "请确认 token 来自飞书知识库 URL（https://xxx.feishu.cn/wiki/NodeToken），"
            "而非普通文档 URL。"
        )

    space_id = node_info["space_id"]
    if not space_id:
        return f"❌ 无法获取知识库 space_id，请检查 token 是否正确。"

    async def _list_children(parent_token: str, depth: int) -> list[dict]:
        """Return flat list of {title, node_token, obj_token, has_child, depth}."""
        async with httpx.AsyncClient(timeout=15) as client:
            resp = await client.get(
                f"https://open.feishu.cn/open-apis/wiki/v2/spaces/{space_id}/nodes",
                headers=headers,
                params={"parent_node_token": parent_token, "page_size": 50},
            )
        data = resp.json()
        if data.get("code") != 0:
            return []
        items = data.get("data", {}).get("items", [])
        result = []
        for item in items:
            entry = {
                "title": item.get("title", "(无标题)"),
                "node_token": item.get("node_token", ""),
                "obj_token": item.get("obj_token", ""),
                "has_child": item.get("has_child", False),
                "depth": depth,
            }
            result.append(entry)
            if recursive and entry["has_child"] and depth < 2:
                children = await _list_children(entry["node_token"], depth + 1)
                result.extend(children)
        return result

    pages = await _list_children(node_token, 0)
    if not pages:
        return f"📂 Wiki 页面 `{node_token}` 下没有子页面。"

    lines = [f"📂 Wiki 页面 `{node_token}` 的子页面（共 {len(pages)} 个）：\n"]
    for p in pages:
        indent = "  " * p["depth"]
        child_hint = " _(有子页面)_" if p["has_child"] else ""
        lines.append(
            f"{indent}• **{p['title']}**{child_hint}\n"
            f"{indent}  node_token: `{p['node_token']}`\n"
            f"{indent}  obj_token: `{p['obj_token']}`"
        )
    lines.append(
        "\n💡 用 `feishu_doc_read(document_token=\"<node_token>\")` 读取每个子页面的内容。"
        "\n   对有子页面的条目，再次调用 `feishu_wiki_list(node_token=\"...\")` 继续展开。"
    )
    return "\n".join(lines)


async def _feishu_doc_read(agent_id: uuid.UUID, arguments: dict) -> str:
    import httpx
    document_token = arguments.get("document_token", "").strip()
    if not document_token:
        return "❌ Missing required argument 'document_token'"
    max_chars = min(int(arguments.get("max_chars", 6000)), 20000)

    creds = await _get_feishu_token(agent_id)
    if not creds:
        return "❌ Agent has no Feishu channel configured."
    _, token = creds

    # Auto-detect wiki node tokens: try get_node first and use obj_token for reading
    read_token = document_token
    wiki_hint = ""
    node_info = await _feishu_wiki_get_node(document_token, token)
    if node_info and node_info.get("obj_token"):
        read_token = node_info["obj_token"]
        if node_info.get("has_child"):
            wiki_hint = (
                "\n\n> 💡 这是一个 Wiki 目录页，它有多个子页面。"
                "使用 `feishu_wiki_list` 工具（传入相同的 node_token）可以查看所有子页面列表。"
            )

    async with httpx.AsyncClient(timeout=20) as client:
        resp = await client.get(
            f"https://open.feishu.cn/open-apis/docx/v1/documents/{read_token}/raw_content",
            headers={"Authorization": f"Bearer {token}"},
            params={"lang": 0},
        )

    data = resp.json()
    if data.get("code") != 0:
        return f"❌ Failed to read document: {data.get('msg')} (code {data.get('code')})"

    content = data.get("data", {}).get("content", "")
    if not content:
        return f"📄 Document '{document_token}' is empty.{wiki_hint}"

    truncated = ""
    if len(content) > max_chars:
        content = content[:max_chars]
        truncated = f"\n\n_(Truncated to {max_chars} chars)_"

    return f"📄 **Document content** (`{document_token}`):\n\n{content}{truncated}{wiki_hint}"


async def _feishu_doc_create(agent_id: uuid.UUID, arguments: dict) -> str:
    import httpx
    title = arguments.get("title", "").strip()
    if not title:
        return "❌ Missing required argument 'title'"

    creds = await _get_feishu_token(agent_id)
    if not creds:
        return "❌ Agent has no Feishu channel configured."
    _, token = creds
    headers = {"Authorization": f"Bearer {token}"}

    body: dict = {"title": title}
    if arguments.get("folder_token"):
        body["folder_token"] = arguments["folder_token"]

    async with httpx.AsyncClient(timeout=20) as client:
        resp = await client.post(
            "https://open.feishu.cn/open-apis/docx/v1/documents",
            json=body,
            headers=headers,
        )

    data = resp.json()
    if data.get("code") != 0:
        return f"❌ Failed to create document: {data.get('msg')} (code {data.get('code')})"

    doc_token = data.get("data", {}).get("document", {}).get("document_id", "")
    doc_url = f"https://bytedance.larkoffice.com/docx/{doc_token}"

    # Auto-share with the Feishu sender so they can access the document
    share_note = ""
    try:
        sender_open_id = channel_feishu_sender_open_id.get(None)
        if sender_open_id and doc_token:
            async with httpx.AsyncClient(timeout=10) as client:
                share_resp = await client.post(
                    f"https://open.feishu.cn/open-apis/drive/v1/permissions/{doc_token}/members",
                    params={"type": "docx", "need_notification": "false"},
                    json={
                        "member_type": "openid",
                        "member_id": sender_open_id,
                        "perm": "full_access",
                    },
                    headers=headers,
                )
            sr = share_resp.json()
            if sr.get("code") == 0:
                share_note = "\n✅ 已自动为你开通访问权限。"
            else:
                share_note = f"\n⚠️ 自动授权失败（{sr.get('code')}），你可能需要手动在飞书前端打开文档。"
    except Exception as _e:
        share_note = f"\n⚠️ 自动授权异常: {_e}"

    return (
        f"✅ 文档创建成功！{share_note}\n"
        f"标题：{title}\n"
        f"Token：{doc_token}\n"
        f"🔗 访问链接：{doc_url}\n"
        f"下一步：调用 feishu_doc_append(document_token=\"{doc_token}\", content=\"...\") 写入正文内容。"
    )


def _parse_inline_markdown(text: str) -> list[dict]:
    """Parse inline markdown (bold, italic, strikethrough) into Feishu text_run elements.
    Note: inline `code` is deliberately NOT rendered as inline_code style because
    Feishu's API rejects inline_code inside heading blocks (field validation error).
    Instead, backtick-wrapped text is returned as plain text.
    Empty text_element_style dicts are intentionally omitted to avoid API validation errors.
    """
    import re as _re

    def _make_run(content: str, style: dict | None = None) -> dict:
        run: dict = {"content": content}
        if style:
            run["text_element_style"] = style
        return {"text_run": run}

    elements = []
    # Only handle **bold**, *italic*, ~~strikethrough~~; backticks become plain text
    pattern = r'(\*\*(.+?)\*\*|\*(.+?)\*|~~(.+?)~~|`(.+?)`)'
    pos = 0
    for m in _re.finditer(pattern, text):
        if m.start() > pos:
            elements.append(_make_run(text[pos:m.start()]))
        raw = m.group(0)
        if raw.startswith("**"):
            elements.append(_make_run(m.group(2), {"bold": True}))
        elif raw.startswith("~~"):
            elements.append(_make_run(m.group(4), {"strikethrough": True}))
        elif raw.startswith("`"):
            # Render as plain text to avoid inline_code validation issues in headings
            elements.append(_make_run(m.group(5)))
        else:
            elements.append(_make_run(m.group(3), {"italic": True}))
        pos = m.end()
    if pos < len(text):
        elements.append(_make_run(text[pos:]))
    if not elements:
        elements.append(_make_run(text or " "))
    return elements


def _markdown_to_feishu_blocks(markdown: str) -> list[dict]:
    """Convert Markdown text to Feishu docx v1 block list.

    Supported:
      # / ## / ### / ####  → heading1-4 (block_type 3-6)
      - / * / + text       → bullet      (block_type 12)
      1. text              → ordered     (block_type 13)
      > text               → quote       (block_type 15)
      --- / ***            → divider     (block_type 22)
      ``` ... ```          → code block  (block_type 14)
      plain text           → text        (block_type 2)
      inline **bold** *italic* `code` ~~strike~~  → text_element_style
    """
    import re as _re

    _HEADING_BLOCK = {1: (3, "heading1"), 2: (4, "heading2"),
                      3: (5, "heading3"), 4: (6, "heading4")}

    def _text_block(bt: int, key: str, line: str) -> dict:
        # Omit "style" entirely to avoid Feishu field validation errors on empty style dicts
        return {
            "block_type": bt,
            key: {"elements": _parse_inline_markdown(line)},
        }

    blocks: list[dict] = []
    lines = markdown.splitlines()
    i = 0
    while i < len(lines):
        line = lines[i]

        # ── Code fence ──────────────────────────────────────────────────────
        if line.strip().startswith("```"):
            lang = line.strip()[3:].strip()
            code_lines = []
            i += 1
            while i < len(lines) and not lines[i].strip().startswith("```"):
                code_lines.append(lines[i])
                i += 1
            blocks.append({
                "block_type": 14,
                "code": {
                    "elements": [{"text_run": {"content": "\n".join(code_lines)}}],
                    "style": {"language": 1 if not lang else
                              {"python": 49, "javascript": 22, "js": 22,
                               "typescript": 56, "ts": 56, "bash": 4, "sh": 4,
                               "sql": 53, "java": 21, "go": 17, "rust": 51,
                               "json": 25, "yaml": 60, "html": 19, "css": 10,
                               }.get(lang.lower(), 1)},
                },
            })
            i += 1
            continue

        # ── Divider ──────────────────────────────────────────────────────────
        if _re.fullmatch(r'[-*_]{3,}', line.strip()):
            # block_type 22 = Divider; no extra fields allowed (empty dict causes validation error)
            blocks.append({"block_type": 22})
            i += 1
            continue

        # ── Headings ─────────────────────────────────────────────────────────
        hm = _re.match(r'^(#{1,4})\s+(.*)', line)
        if hm:
            level = min(len(hm.group(1)), 4)
            bt, key = _HEADING_BLOCK[level]
            blocks.append(_text_block(bt, key, hm.group(2)))
            i += 1
            continue

        # ── Bullet list ──────────────────────────────────────────────────────
        if _re.match(r'^[\-\*\+]\s+', line):
            text = _re.sub(r'^[\-\*\+]\s+', '', line)
            blocks.append(_text_block(12, "bullet", text))
            i += 1
            continue

        # ── Ordered list ─────────────────────────────────────────────────────
        if _re.match(r'^\d+\.\s+', line):
            text = _re.sub(r'^\d+\.\s+', '', line)
            blocks.append(_text_block(13, "ordered", text))
            i += 1
            continue

        # ── Blockquote ───────────────────────────────────────────────────────
        if line.startswith("> "):
            blocks.append(_text_block(15, "quote", line[2:]))
            i += 1
            continue

        # ── Empty line → empty text block ────────────────────────────────────
        if line.strip() == "":
            blocks.append({
                "block_type": 2,
                "text": {"elements": [{"text_run": {"content": " "}}]},
            })
            i += 1
            continue

        # ── Markdown table separator line (|---|---| ) → skip ───────────────
        if _re.match(r'^\|[\s\-:]+(\|[\s\-:]+)*\|?\s*$', line.strip()):
            i += 1
            continue

        # ── Markdown table row → plain text ──────────────────────────────────
        if line.strip().startswith("|") and line.strip().endswith("|"):
            # Strip pipe separators and render each cell as plain text
            cells = [c.strip() for c in line.strip().strip("|").split("|")]
            cell_text = "  |  ".join(c for c in cells if c)
            blocks.append(_text_block(2, "text", cell_text))
            i += 1
            continue

        # ── Plain text (with inline formatting) ──────────────────────────────
        blocks.append(_text_block(2, "text", line))
        i += 1

    return blocks


async def _feishu_doc_append(agent_id: uuid.UUID, arguments: dict) -> str:
    import httpx
    document_token = arguments.get("document_token", "").strip()
    content = arguments.get("content", "").strip()
    if not document_token:
        return "❌ Missing required argument 'document_token'"
    if not content:
        return "❌ Missing required argument 'content'"

    creds = await _get_feishu_token(agent_id)
    if not creds:
        return "❌ Agent has no Feishu channel configured."
    _, token = creds
    headers = {"Authorization": f"Bearer {token}"}

    # For wiki node tokens, use the obj_token for the docx API
    node_info = await _feishu_wiki_get_node(document_token, token)
    docx_token = node_info["obj_token"] if (node_info and node_info.get("obj_token")) else document_token

    async with httpx.AsyncClient(timeout=20) as client:
        meta = (await client.get(
            f"https://open.feishu.cn/open-apis/docx/v1/documents/{docx_token}",
            headers=headers,
        )).json()
        if meta.get("code") != 0:
            return f"❌ Cannot access document: {meta.get('msg')}"

        body_block_id = (
            meta.get("data", {}).get("document", {}).get("body", {}).get("block_id")
            or docx_token
        )

        children = _markdown_to_feishu_blocks(content)

        result = (await client.post(
            f"https://open.feishu.cn/open-apis/docx/v1/documents/{docx_token}/blocks/{body_block_id}/children",
            json={"children": children, "index": -1},
            headers=headers,
        )).json()

    if result.get("code") != 0:
        return f"❌ Failed to append: {result.get('msg')} (code {result.get('code')})"

    doc_url = f"https://bytedance.larkoffice.com/docx/{docx_token}"
    return (
        f"✅ 已写入 {len(children)} 个段落到文档。\n"
        f"🔗 文档直链（原文发给用户，勿修改）：{doc_url}"
    )


# ─── Feishu Document Share ────────────────────────────────────────────────────

async def _feishu_doc_share(agent_id: uuid.UUID, arguments: dict) -> str:
    """Manage Feishu document collaborators.
    Automatically handles both regular docx documents (Drive permissions API)
    and Wiki node documents (Wiki space members API).
    """
    import httpx
    import re as _re

    document_token = (arguments.get("document_token") or "").strip()
    action = (arguments.get("action") or "list").strip()
    permission = (arguments.get("permission") or "edit").strip()

    if not document_token:
        return "❌ Missing required argument 'document_token'"

    creds = await _get_feishu_token(agent_id)
    if not creds:
        return "❌ Agent has no Feishu channel configured."
    _, token = creds
    headers = {"Authorization": f"Bearer {token}"}

    # ── Detect if this is a Wiki node token ─────────────────────────────────
    node_info = await _feishu_wiki_get_node(document_token, token)
    is_wiki = node_info is not None
    space_id = node_info.get("space_id", "") if node_info else ""
    obj_token = node_info.get("obj_token", "") if node_info else ""

    # Permission level mapping: Feishu API uses "view" / "edit" / "full_access"
    api_perm = {"view": "view", "edit": "edit", "full_access": "full_access"}.get(permission, "edit")
    # Wiki space role mapping: only "admin" / "member" are valid roles
    wiki_role = "admin" if api_perm in ("edit", "full_access") else "member"

    # ── LIST collaborators ────────────────────────────────────────────────────
    if action == "list":
        use_token = obj_token if (is_wiki and obj_token) else document_token
        async with httpx.AsyncClient(timeout=15) as client:
            resp = await client.get(
                f"https://open.feishu.cn/open-apis/drive/v1/permissions/{use_token}/members",
                params={"type": "docx"},
                headers=headers,
            )
        data = resp.json()
        if data.get("code") != 0:
            _c = data.get("code")
            if _c == 1063003 and is_wiki:
                return (
                    f"ℹ️ 文档 `{document_token}` 是知识库页面，其权限由知识库空间统一管理。\n"
                    "知识库空间 ID：`" + space_id + "`\n"
                    "请直接在飞书知识库中管理成员权限。"
                )
            if _c in (99991672, 99991668):
                return (
                    f"❌ 权限不足（code {_c}）\n"
                    "需要在飞书开放平台开通：\n"
                    "• drive:drive（云文档权限管理）"
                )
            return f"❌ 获取协作者列表失败：{data.get('msg')} (code {_c})"

        members = data.get("data", {}).get("items", [])
        if not members:
            return f"📄 文档 `{document_token}` 当前没有其他协作者。"

        lines = [f"📄 文档 `{document_token}` 的协作者列表（共 {len(members)} 人）：\n"]
        for m in members:
            perm = m.get("perm", "")
            member_type = m.get("member_type", "")
            member_id = m.get("member_id", "")
            _type_label = {"openid": "用户", "openchat": "群组", "opendepartmentid": "部门"}.get(member_type, member_type)
            lines.append(f"• {_type_label} `{member_id}` | 权限: **{perm}**")
        return "\n".join(lines)

    # ── ADD / REMOVE collaborators ─────────────────────────────────────────────
    member_names: list[str] = list(arguments.get("member_names") or [])
    member_open_ids: list[str] = list(arguments.get("member_open_ids") or [])

    if not member_names and not member_open_ids:
        return "❌ 请提供 member_names（姓名列表）或 member_open_ids（open_id 列表）"

    # Resolve names → open_ids
    resolved: list[tuple[str, str]] = []  # (display_name, open_id)
    for name in member_names:
        sr = await _feishu_user_search(agent_id, {"name": name})
        m = _re.search(r'open_id: `(ou_[A-Za-z0-9]+)`', sr)
        if m:
            resolved.append((name, m.group(1)))
        else:
            resolved.append((name, ""))

    for oid in member_open_ids:
        if oid:
            resolved.append((oid, oid))

    results = []
    async with httpx.AsyncClient(timeout=15) as client:
        for display, oid in resolved:
            if not oid:
                results.append(f"❌ 无法找到「{display}」的 open_id，跳过")
                continue

            if action == "add":
                # ── Wiki node: use wiki space members API ──────────────────
                if is_wiki and space_id:
                    resp = await client.post(
                        f"https://open.feishu.cn/open-apis/wiki/v2/spaces/{space_id}/members",
                        json={"member_type": "openid", "member_id": oid, "member_role": wiki_role},
                        headers=headers,
                    )
                    d = resp.json()
                    _c = d.get("code")
                    if _c == 0:
                        results.append(f"✅ 已将「{display}」加入知识库空间（角色：{wiki_role}）")
                    elif _c == 131008:
                        results.append(f"ℹ️ 「{display}」已经是知识库成员，无需重复添加")
                    elif _c == 131101:
                        # Public wiki space — everyone already has access
                        results.append(
                            f"ℹ️ 这是一个**公开知识库**，所有人已可访问。\n"
                            f"「{display}」无需单独添加权限。"
                        )
                    else:
                        results.append(f"❌ 添加「{display}」到知识库失败：{d.get('msg')} (code {_c})")
                    continue

                # ── Regular docx: use Drive permissions API ────────────────
                body = {
                    "member_type": "openid",
                    "member_id": oid,
                    "perm": api_perm,
                }
                resp = await client.post(
                    f"https://open.feishu.cn/open-apis/drive/v1/permissions/{document_token}/members",
                    json=body,
                    headers=headers,
                    params={"type": "docx"},
                )
                d = resp.json()
                if d.get("code") == 0:
                    results.append(f"✅ 已将「{display}」添加为**{permission}**权限协作者")
                else:
                    _c = d.get("code")
                    if _c == 99992402:
                        # Feishu platform policy: you cannot add yourself as a collaborator via API.
                        # Permissions must be granted by others, or set manually in the UI.
                        results.append(
                            f"⚠️ 飞书平台安全限制：无法通过 API 为自己添加协作权限。\n"
                            f"请手动操作：打开文档 → 右上角「分享」→ 添加自己并设置权限。"
                        )
                    elif _c in (99991672, 99991668):
                        return (
                            f"❌ 权限不足（code {_c}）\n"
                            "需要在飞书开放平台开通：\n"
                            "• drive:drive（云文档权限管理）"
                        )
                    else:
                        results.append(f"❌ 添加「{display}」失败：{d.get('msg')} (code {_c})")

            elif action == "remove":
                if is_wiki and space_id:
                    resp = await client.delete(
                        f"https://open.feishu.cn/open-apis/wiki/v2/spaces/{space_id}/members/{oid}",
                        headers=headers,
                        params={"member_type": "openid"},
                    )
                    d = resp.json()
                    if d.get("code") == 0:
                        results.append(f"✅ 已将「{display}」从知识库移除")
                    else:
                        results.append(f"❌ 移除「{display}」失败：{d.get('msg')} (code {d.get('code')})")
                    continue

                resp = await client.delete(
                    f"https://open.feishu.cn/open-apis/drive/v1/permissions/{document_token}/members/{oid}",
                    headers=headers,
                    params={"type": "docx", "member_type": "openid"},
                )
                d = resp.json()
                if d.get("code") == 0:
                    results.append(f"✅ 已移除「{display}」的协作权限")
                else:
                    results.append(f"❌ 移除「{display}」失败：{d.get('msg')} (code {d.get('code')})")

    return "\n".join(results) if results else "没有需要处理的成员"


# ─── Feishu Calendar Tools ────────────────────────────────────────────────────

async def _feishu_calendar_list(agent_id: uuid.UUID, arguments: dict) -> str:
    import httpx
    import re as _re
    from datetime import timedelta as _td

    user_email = arguments.get("user_email", "").strip()

    creds = await _get_feishu_token(agent_id)
    if not creds:
        return "❌ Agent has no Feishu channel configured."
    _, token = creds

    now = datetime.now(timezone.utc)

    def _to_iso(t: str | None, default: datetime) -> str:
        """Return an ISO-8601 string with timezone for freebusy API."""
        if not t:
            return default.strftime("%Y-%m-%dT%H:%M:%S+00:00")
        if _re.fullmatch(r'\d+', t.strip()):
            from datetime import datetime as _dt2
            return _dt2.fromtimestamp(int(t.strip()), tz=timezone.utc).strftime("%Y-%m-%dT%H:%M:%S+00:00")
        return t.strip()

    def _to_unix(t: str | None, default: datetime) -> str:
        """Convert ISO-8601 / Unix string / None to Unix timestamp string."""
        if not t:
            return str(int(default.timestamp()))
        if _re.fullmatch(r'\d+', t.strip()):
            return t.strip()
        try:
            from datetime import datetime as _dt2
            for fmt in ("%Y-%m-%dT%H:%M:%S%z", "%Y-%m-%dT%H:%M:%SZ", "%Y-%m-%dT%H:%M:%S"):
                try:
                    dt = _dt2.strptime(t.strip(), fmt)
                    if dt.tzinfo is None:
                        dt = dt.replace(tzinfo=timezone.utc)
                    return str(int(dt.timestamp()))
                except ValueError:
                    continue
            from dateutil import parser as _dp
            return str(int(_dp.parse(t).timestamp()))
        except Exception:
            return str(int(default.timestamp()))

    start_arg = arguments.get("start_time")
    end_arg = arguments.get("end_time")
    start_ts = _to_unix(start_arg, now)
    end_ts = _to_unix(end_arg, now + _td(days=7))
    start_iso = _to_iso(start_arg, now)
    end_iso = _to_iso(end_arg, now + _td(days=7))

    # ── 1. Query sender's real freebusy from Feishu Calendar ─────────────────
    sender_open_id = channel_feishu_sender_open_id.get(None)
    # Allow explicit override via argument
    if arguments.get("user_open_id"):
        sender_open_id = arguments["user_open_id"]
    elif user_email:
        resolved = await _feishu_resolve_open_id(token, user_email)
        if resolved:
            sender_open_id = resolved

    freebusy_section = ""
    if sender_open_id:
        try:
            async with httpx.AsyncClient(timeout=10) as fb_client:
                fb_resp = await fb_client.post(
                    "https://open.feishu.cn/open-apis/calendar/v4/freebusy/list",
                    headers={"Authorization": f"Bearer {token}"},
                    params={"user_id_type": "open_id"},
                    json={
                        "time_min": start_iso,
                        "time_max": end_iso,
                        "user_id": sender_open_id,
                    },
                )
            fb_data = fb_resp.json()
            if fb_data.get("code") == 0:
                busy_slots = fb_data.get("data", {}).get("freebusy_list", [])
                if busy_slots:
                    from datetime import datetime as _dt2
                    from zoneinfo import ZoneInfo
                    tz_cn = ZoneInfo("Asia/Shanghai")
                    busy_lines = []
                    for slot in sorted(busy_slots, key=lambda x: x.get("start_time", "")):
                        try:
                            s = _dt2.fromisoformat(slot["start_time"]).astimezone(tz_cn).strftime("%H:%M")
                            e = _dt2.fromisoformat(slot["end_time"]).astimezone(tz_cn).strftime("%H:%M")
                            busy_lines.append(f"  🔴 {s}–{e}")
                        except Exception:
                            busy_lines.append(f"  🔴 {slot.get('start_time')}–{slot.get('end_time')}")
                    freebusy_section = f"\n📌 **用户真实日历（忙碌时段）**：\n" + "\n".join(busy_lines)
                else:
                    freebusy_section = "\n📌 **用户真实日历**：该时段全部空闲。"
        except Exception as _fe:
            freebusy_section = f"\n⚠️ Freebusy 查询异常: {_fe}"

    # ── 2. Also list bot's own calendar events ───────────────────────────────
    agent_cal_id, cal_err = await _get_agent_calendar_id(token)
    if not agent_cal_id:
        # Return freebusy results even if bot calendar fails
        if freebusy_section:
            return freebusy_section.strip()
        return cal_err or "❌ Failed to retrieve agent's primary calendar ID."

    # Note: page_size is NOT a valid param for this API — omit it entirely
    params: dict = {}
    if start_ts:
        params["start_time"] = start_ts
    if end_ts:
        params["end_time"] = end_ts

    async with httpx.AsyncClient(timeout=20) as client:
        resp = await client.get(
            f"https://open.feishu.cn/open-apis/calendar/v4/calendars/{agent_cal_id}/events",
            headers={"Authorization": f"Bearer {token}"},
            params=params,
        )

    data = resp.json()
    if data.get("code") != 0:
        if freebusy_section:
            return freebusy_section.strip()
        return f"❌ Calendar API error: {data.get('msg')} (code {data.get('code')})"

    items = data.get("data", {}).get("items", [])
    if not items and not freebusy_section:
        return "📅 该时间段内没有日程。"

    lines = []
    if items:
        lines.append(f"📅 Bot 日历共 {len(items)} 个日程：\n")
    for ev in items:
        summary = ev.get("summary", "(no title)")
        start = ev.get("start_time", {}).get("timestamp", "")
        end_t = ev.get("end_time", {}).get("timestamp", "")
        location = ev.get("location", {}).get("name", "")
        event_id = ev.get("event_id", "")
        try:
            from datetime import datetime as _dt
            s = _dt.fromtimestamp(int(start), tz=timezone.utc).strftime("%m-%d %H:%M") if start else "?"
            e = _dt.fromtimestamp(int(end_t), tz=timezone.utc).strftime("%H:%M") if end_t else "?"
        except Exception:
            s, e = start, end_t
        loc_str = f" | 📍{location}" if location else ""
        lines.append(f"- **{summary}** | 🕐{s}–{e}{loc_str}  (ID: `{event_id}`)")

    if freebusy_section:
        lines.append(freebusy_section)

    return "\n".join(lines) if lines else "📅 该时间段内没有日程。"


async def _feishu_calendar_create(agent_id: uuid.UUID, arguments: dict) -> str:
    import httpx

    user_email = arguments.get("user_email", "").strip()
    summary = arguments.get("summary", "").strip()
    start_time = arguments.get("start_time", "").strip()
    end_time = arguments.get("end_time", "").strip()

    for f, v in [("summary", summary), ("start_time", start_time), ("end_time", end_time)]:
        if not v:
            return f"❌ Missing required argument '{f}'"

    creds = await _get_feishu_token(agent_id)
    if not creds:
        return "❌ Agent has no Feishu channel configured."
    _, token = creds

    # Resolve organizer open_id from email — soft failure
    organizer_open_id: str | None = None
    if user_email:
        organizer_open_id = await _feishu_resolve_open_id(token, user_email)
        if not organizer_open_id:
            print(f"[Feishu Calendar] Could not resolve open_id for '{user_email}', continuing without organizer invite")

    agent_cal_id, cal_err = await _get_agent_calendar_id(token)
    if not agent_cal_id:
        return cal_err or "❌ Failed to retrieve agent's primary calendar ID."

    tz = arguments.get("timezone", "Asia/Shanghai")
    body: dict = {
        "summary": summary,
        "start_time": {"timestamp": str(int(_iso_to_ts(start_time))), "timezone": tz},
        "end_time": {"timestamp": str(int(_iso_to_ts(end_time))), "timezone": tz},
    }
    if arguments.get("description"):
        body["description"] = arguments["description"]
    if arguments.get("location"):
        body["location"] = {"name": arguments["location"]}

    async with httpx.AsyncClient(timeout=20) as client:
        resp = await client.post(
            f"https://open.feishu.cn/open-apis/calendar/v4/calendars/{agent_cal_id}/events",
            json=body,
            headers={"Authorization": f"Bearer {token}"},
        )

    data = resp.json()
    if data.get("code") != 0:
        return f"❌ Failed to create event: {data.get('msg')} (code {data.get('code')})"

    event_id = data.get("data", {}).get("event", {}).get("event_id", "")

    # Collect all attendee open_ids to invite
    attendee_open_ids: list[str] = []
    attendee_display: list[str] = []  # for summary message

    # 1. Direct open_ids provided by caller
    for oid in (arguments.get("attendee_open_ids") or []):
        if oid and oid not in attendee_open_ids:
            attendee_open_ids.append(oid)
            attendee_display.append(oid)

    # 2. Names → look up via feishu_user_search
    import re as _re_oid
    for aname in (arguments.get("attendee_names") or []):
        aname = aname.strip()
        if not aname:
            continue
        _sr = await _feishu_user_search(agent_id, {"name": aname})
        _m = _re_oid.search(r'open_id: `(ou_[A-Za-z0-9]+)`', _sr)
        if _m:
            _oid = _m.group(1)
            if _oid not in attendee_open_ids:
                attendee_open_ids.append(_oid)
                attendee_display.append(aname)
        else:
            print(f"[Calendar] Could not resolve attendee '{aname}': {_sr[:100]}")

    # 3. From explicit attendee_emails
    attendee_emails: list[str] = list(arguments.get("attendee_emails") or [])
    if user_email and user_email not in attendee_emails:
        attendee_emails.append(user_email)
    for email in attendee_emails[:20]:
        oid = await _feishu_resolve_open_id(token, email)
        if oid and oid not in attendee_open_ids:
            attendee_open_ids.append(oid)
            attendee_display.append(email)

    # 4. Auto-invite the Feishu message sender (from context var)
    sender_oid = channel_feishu_sender_open_id.get(None)
    if sender_oid and sender_oid not in attendee_open_ids:
        attendee_open_ids.append(sender_oid)

    if attendee_open_ids and event_id:
        async with httpx.AsyncClient(timeout=20) as client:
            for oid in attendee_open_ids:
                await client.post(
                    f"https://open.feishu.cn/open-apis/calendar/v4/calendars/{agent_cal_id}/events/{event_id}/attendees",
                    json={"attendees": [{"type": "user", "user_id": oid}]},
                    headers={"Authorization": f"Bearer {token}"},
                    params={"user_id_type": "open_id"},
                )

    att_str = f"\n**参与人**: {', '.join(attendee_display)}" if attendee_display else ""
    invite_note = "\n（已向您发送日历邀请，请在飞书日历中确认）" if attendee_open_ids else ""
    return (
        f"✅ 日历事件已创建！\n"
        f"**标题**: {summary}\n"
        f"**时间**: {start_time} → {end_time}{att_str}\n"
        f"**Event ID**: `{event_id}`{invite_note}"
    )


async def _feishu_calendar_update(agent_id: uuid.UUID, arguments: dict) -> str:
    import httpx

    user_email = arguments.get("user_email", "").strip()
    event_id = arguments.get("event_id", "").strip()
    if not user_email or not event_id:
        return "❌ Both 'user_email' and 'event_id' are required."

    creds = await _get_feishu_token(agent_id)
    if not creds:
        return "❌ Agent has no Feishu channel configured."
    _, token = creds

    open_id = await _feishu_resolve_open_id(token, user_email)
    if not open_id:
        return f"❌ User '{user_email}' not found."

    agent_cal_id, cal_err = await _get_agent_calendar_id(token)
    if not agent_cal_id:
        return cal_err or "❌ Failed to retrieve agent's primary calendar ID."

    patch: dict = {}
    tz = arguments.get("timezone", "Asia/Shanghai")
    if arguments.get("summary"):
        patch["summary"] = arguments["summary"]
    if arguments.get("description"):
        patch["description"] = arguments["description"]
    if arguments.get("location"):
        patch["location"] = {"name": arguments["location"]}
    if arguments.get("start_time"):
        patch["start_time"] = {"timestamp": str(int(_iso_to_ts(arguments["start_time"]))), "timezone": tz}
    if arguments.get("end_time"):
        patch["end_time"] = {"timestamp": str(int(_iso_to_ts(arguments["end_time"]))), "timezone": tz}

    if not patch:
        return "ℹ️ No fields to update."

    async with httpx.AsyncClient(timeout=20) as client:
        resp = await client.patch(
            f"https://open.feishu.cn/open-apis/calendar/v4/calendars/{agent_cal_id}/events/{event_id}",
            json=patch,
            headers={"Authorization": f"Bearer {token}"},
        )

    data = resp.json()
    if data.get("code") != 0:
        return f"❌ Failed to update: {data.get('msg')} (code {data.get('code')})"

    return f"✅ Event `{event_id}` updated. Changed: {', '.join(patch.keys())}."


async def _feishu_calendar_delete(agent_id: uuid.UUID, arguments: dict) -> str:
    import httpx

    user_email = arguments.get("user_email", "").strip()
    event_id = arguments.get("event_id", "").strip()
    if not user_email or not event_id:
        return "❌ Both 'user_email' and 'event_id' are required."

    creds = await _get_feishu_token(agent_id)
    if not creds:
        return "❌ Agent has no Feishu channel configured."
    _, token = creds

    open_id = await _feishu_resolve_open_id(token, user_email)
    if not open_id:
        return f"❌ User '{user_email}' not found."

    agent_cal_id, cal_err = await _get_agent_calendar_id(token)
    if not agent_cal_id:
        return cal_err or "❌ Failed to retrieve agent's primary calendar ID."

    async with httpx.AsyncClient(timeout=20) as client:
        resp = await client.delete(
            f"https://open.feishu.cn/open-apis/calendar/v4/calendars/{agent_cal_id}/events/{event_id}",
            headers={"Authorization": f"Bearer {token}"},
        )

    data = resp.json()
    if data.get("code") != 0:
        return f"❌ Failed to delete: {data.get('msg')} (code {data.get('code')})"

    return f"✅ Event `{event_id}` deleted successfully."


# ─── Feishu User Search ───────────────────────────────────────────────────────

async def _feishu_user_search(agent_id: uuid.UUID, arguments: dict) -> str:
    """Search for colleagues in the Feishu directory by name.

    Strategy:
    1. Search local contacts cache (populated when anyone messages the bot).
    2. Fall back to Contact v3 GET /users/{open_id} if we find a match by email.
    The cache is populated by feishu.py each time a message sender is resolved.
    """
    import httpx
    import json as _json
    import pathlib as _pl

    name = (arguments.get("name") or "").strip()
    if not name:
        return "❌ Missing required argument 'name'"

    creds = await _get_feishu_token(agent_id)
    if not creds:
        return "❌ Agent has no Feishu channel configured."
    _, token = creds

    # ── Load local contacts cache ─────────────────────────────────────────────
    _cache_file = _pl.Path(f"/data/workspaces/{agent_id}/feishu_contacts_cache.json")
    _cached_users: list[dict] = []
    try:
        if _cache_file.exists():
            _raw = _json.loads(_cache_file.read_text())
            _cached_users = _raw.get("users", [])
    except Exception:
        pass

    name_lower = name.lower()

    def _matches(u: dict) -> bool:
        return (
            name_lower in (u.get("name") or "").lower()
            or name_lower in (u.get("en_name") or "").lower()
        )

    matched = [u for u in _cached_users if _matches(u)]

    if matched:
        lines = [f"🔍 找到 {len(matched)} 位匹配「{name}」的用户：\n"]
        for u in matched:
            open_id = u.get("open_id", "")
            display_name = u.get("name", "")
            en_name = u.get("en_name", "")
            email = u.get("email", "")
            lines.append(f"• **{display_name}**{'（' + en_name + '）' if en_name else ''}")
            lines.append(f"  open_id: `{open_id}`")
            if email:
                lines.append(f"  邮箱: {email}")
        return "\n".join(lines)

    # ── Cache miss: try resolving via email (if caller provided one) ──────────
    # Also try Contact v3 with a known open_id from DB users
    try:
        from app.database import async_session as _async_session
        from sqlalchemy import select as _sa_select
        from app.models.user import User as _User
        async with _async_session() as _db:
            _r = await _db.execute(
                _sa_select(_User).where(_User.display_name.ilike(f"%{name}%"))
            )
            _platform_users = _r.scalars().all()
        for _pu in _platform_users:
            _oid = getattr(_pu, "feishu_open_id", None)
            _email = getattr(_pu, "email", None)
            if _oid:
                async with httpx.AsyncClient(timeout=10) as _c:
                    _r2 = await _c.get(
                        f"https://open.feishu.cn/open-apis/contact/v3/users/{_oid}",
                        params={"user_id_type": "open_id"},
                        headers={"Authorization": f"Bearer {token}"},
                    )
                _d2 = _r2.json()
                if _d2.get("code") == 0:
                    _ui = _d2.get("data", {}).get("user", {})
                    _found_name = _ui.get("name", _pu.display_name or "")
                    _found_email = _ui.get("email", "") or _email or ""
                    return (
                        f"🔍 找到匹配「{name}」的用户：\n\n"
                        f"• **{_found_name}**\n"
                        f"  open_id: `{_oid}`\n"
                        + (f"  邮箱: {_found_email}\n" if _found_email else "")
                    )
    except Exception:
        pass

    total = len(_cached_users)
    if total == 0:
        return (
            f"❌ 本地通讯录缓存为空，暂时无法搜索「{name}」。\n\n"
            "通讯录缓存会在同事向机器人发消息时自动建立。\n"
            "如果「覃睿」从未给机器人发过消息，可以请他先给机器人发一条消息，"
            "之后就能直接搜索到他了。\n\n"
            "或者，请直接告诉我「覃睿」的飞书 open_id 或邮箱，我可以立刻操作。"
        )
    return (
        f"❌ 未在本地通讯录（已缓存 {total} 人）中找到「{name}」。\n\n"
        "通讯录缓存来自给机器人发过消息的同事。\n"
        "如果「{name}」从未给机器人发消息，请他先发一条，之后即可自动识别。\n"
        "或者请直接提供其飞书 open_id / 工作邮箱。"
    )


async def _feishu_contacts_refresh(agent_id: uuid.UUID) -> None:
    """Force-clear the local contacts cache so next search re-fetches from API."""
    import pathlib as _pl
    _cache_file = _pl.Path("/data/workspaces") / str(agent_id) / "feishu_contacts_cache.json"
    try:
        if _cache_file.exists():
            _cache_file.unlink()
    except Exception:
        pass


# ─── Email Tool Helpers ─────────────────────────────────────

async def _get_email_config(agent_id: uuid.UUID) -> dict:
    """Retrieve per-agent email config from the send_email tool's AgentTool config."""
    from app.models.tool import Tool, AgentTool

    async with async_session() as db:
        # Find the send_email tool
        r = await db.execute(select(Tool).where(Tool.name == "send_email"))
        tool = r.scalar_one_or_none()
        if not tool:
            return {}

        # Get per-agent config
        at_r = await db.execute(
            select(AgentTool).where(
                AgentTool.agent_id == agent_id,
                AgentTool.tool_id == tool.id,
            )
        )
        at = at_r.scalar_one_or_none()
        agent_config = (at.config or {}) if at else {}
        # Merge global + agent override
        return {**(tool.config or {}), **agent_config}


async def _handle_email_tool(tool_name: str, agent_id: uuid.UUID, ws: Path, arguments: dict) -> str:
    """Dispatch email tool calls to the email_service module."""
    from app.services.email_service import send_email, read_emails, reply_email

    config = await _get_email_config(agent_id)
    if not config.get("email_address") or not config.get("auth_code"):
        return (
            "❌ Email not configured for this agent.\n\n"
            "Please go to Agent → Tools → Send Email → Config to set up your email:\n"
            "1. Select your email provider\n"
            "2. Enter your email address\n"
            "3. Enter your authorization code (not your login password)"
        )

    try:
        if tool_name == "send_email":
            return await send_email(
                config=config,
                to=arguments.get("to", ""),
                subject=arguments.get("subject", ""),
                body=arguments.get("body", ""),
                cc=arguments.get("cc"),
                attachments=arguments.get("attachments"),
                workspace_path=ws,
            )
        elif tool_name == "read_emails":
            return await read_emails(
                config=config,
                limit=arguments.get("limit", 10),
                search=arguments.get("search"),
                folder=arguments.get("folder", "INBOX"),
            )
        elif tool_name == "reply_email":
            return await reply_email(
                config=config,
                message_id=arguments.get("message_id", ""),
                body=arguments.get("body", ""),
            )
        else:
            return f"❌ Unknown email tool: {tool_name}"
    except Exception as e:
        return f"❌ Email tool error: {str(e)[:200]}"

